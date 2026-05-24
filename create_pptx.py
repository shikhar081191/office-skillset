"""Create .pptx files using python-pptx (pure Python, no npm needed).

This replaces the pptxgenjs / npm workflow for PowerPoint creation.
Includes design helpers for colour theming, layouts, and visual elements.

Usage:
    python create_pptx.py          # runs demo
    from create_pptx import PptxBuilder
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE


# ── colour palettes (from skill) ─────────────────────────────────────────────

DEFAULT_ROLES = {
    "background": "FFFFFF",
    "surface": "F5F5F5",
    "surface_alt": "EEEEEE",
    "border": "D9D9D9",
    "muted_text": "666666",
    "positive": "2E7D32",
    "warning": "D28B00",
    "negative": "C33D3D",
    "highlight": "FFF3CC",
}

PALETTES = {
    "blackrock": {
        "primary":    "000000",
        "secondary":  "FF4713",
        "accent":     "FFCE00",
        "text_dark":  "000000",
        "text_light": "FFFFFF",
    },
    "midnight_executive": {
        "primary":    "1E2761",
        "secondary":  "CADCFC",
        "accent":     "FFFFFF",
        "text_dark":  "1E2761",
        "text_light": "FFFFFF",
    },
    "coral_energy": {
        "primary":    "F96167",
        "secondary":  "F9E795",
        "accent":     "2F3C7E",
        "text_dark":  "2F3C7E",
        "text_light": "FFFFFF",
    },
    "teal_trust": {
        "primary":    "028090",
        "secondary":  "00A896",
        "accent":     "02C39A",
        "text_dark":  "065A82",
        "text_light": "FFFFFF",
    },
    "charcoal_minimal": {
        "primary":    "36454F",
        "secondary":  "F2F2F2",
        "accent":     "212121",
        "text_dark":  "212121",
        "text_light": "FFFFFF",
    },
    "warm_terracotta": {
        "primary":    "B85042",
        "secondary":  "E7E8D1",
        "accent":     "A7BEAE",
        "text_dark":  "B85042",
        "text_light": "FFFFFF",
    },
}


def resolve_palette(palette=None) -> dict:
    """Return a complete palette from a name or an external role dictionary."""
    selected = palette if palette is not None else "blackrock"
    if isinstance(selected, str):
        selected = PALETTES.get(selected, PALETTES["blackrock"])
    resolved = dict(DEFAULT_ROLES)
    resolved.update(selected)
    required = {"primary", "secondary", "accent", "text_dark", "text_light"}
    missing = required.difference(resolved)
    if missing:
        raise ValueError(f"Palette missing required keys: {', '.join(sorted(missing))}")
    return resolved


def register_palette(name: str, palette: dict) -> None:
    """Register a reusable named palette supplied by a team or project."""
    PALETTES[name] = resolve_palette(palette)


def hex_to_rgb(hex_color: str) -> RGBColor:
    h = hex_color.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


class PptxBuilder:
    """
    Fluent builder for PowerPoint presentations.

    palette: one of the keys in PALETTES, or a dict with the core brand keys.
             Optional semantic roles customize surfaces and data colours.
    """

    SLIDE_W = Inches(13.33)
    SLIDE_H = Inches(7.5)

    def __init__(self, palette="blackrock"):
        self.prs = Presentation()
        self.prs.slide_width = self.SLIDE_W
        self.prs.slide_height = self.SLIDE_H

        self.palette = resolve_palette(palette)

        self._blank_layout = self.prs.slide_layouts[6]  # blank

    # ── slide types ──────────────────────────────────────────────────────────

    def title_slide(self, title: str, subtitle: str = "") -> "PptxBuilder":
        """Full-bleed dark title slide."""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._fill_background(slide, self.palette["primary"])

        # Title
        self._add_text_box(
            slide, title,
            x=Inches(0.8), y=Inches(2.5), w=Inches(11.7), h=Inches(1.5),
            font_size=44, bold=True,
            color=self.palette["text_light"],
            align=PP_ALIGN.LEFT,
        )

        # Subtitle
        if subtitle:
            self._add_text_box(
                slide, subtitle,
                x=Inches(0.8), y=Inches(4.2), w=Inches(11.7), h=Inches(0.8),
                font_size=20, bold=False,
                color=self.palette["secondary"],
                align=PP_ALIGN.LEFT,
            )

        # Accent bar (bottom strip)
        self._add_rectangle(
            slide,
            x=0, y=Inches(6.9), w=self.SLIDE_W, h=Inches(0.6),
            fill_color=self.palette["accent"],
        )
        return self

    def content_slide(self, title: str, body: str | list,
                      dark_bg=False) -> "PptxBuilder":
        """Standard title + body text slide."""
        slide = self.prs.slides.add_slide(self._blank_layout)
        bg_color = self.palette["primary"] if dark_bg else self.palette["background"]
        text_color = self.palette["text_light"] if dark_bg else self.palette["text_dark"]
        self._fill_background(slide, bg_color)

        # Title bar
        self._add_rectangle(
            slide,
            x=0, y=0, w=self.SLIDE_W, h=Inches(1.2),
            fill_color=self.palette["primary"],
        )
        self._add_text_box(
            slide, title,
            x=Inches(0.5), y=Inches(0.1), w=Inches(12.3), h=Inches(1.0),
            font_size=28, bold=True, color=self.palette["text_light"],
        )

        # Body
        if isinstance(body, list):
            body_text = "\n".join(f"•  {item}" for item in body)
            # Scale font up when bullet count is low so canvas fills visually
            font_size = 20 if len(body) <= 3 else 16 if len(body) <= 5 else 14
        else:
            body_text = body
            font_size = 18 if len(body.split()) <= 40 else 14

        self._add_text_box(
            slide, body_text,
            x=Inches(0.8), y=Inches(1.5), w=Inches(11.7), h=Inches(5.5),
            font_size=font_size, bold=False, color=text_color,
            align=PP_ALIGN.LEFT, wrap=True, auto_fit=True,
        )
        return self

    def two_column_slide(self, title: str,
                         left_heading: str, left_body: str | list,
                         right_heading: str, right_body: str | list) -> "PptxBuilder":
        """Two-column content slide."""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._fill_background(slide, self.palette["background"])

        # Title bar
        self._add_rectangle(slide, x=0, y=0, w=self.SLIDE_W, h=Inches(1.2),
                             fill_color=self.palette["primary"])
        self._add_text_box(slide, title,
                           x=Inches(0.5), y=Inches(0.1), w=Inches(12.3), h=Inches(1.0),
                           font_size=28, bold=True, color=self.palette["text_light"])

        col_w = Inches(5.8)
        for x_offset, heading, body in [
            (Inches(0.5), left_heading, left_body),
            (Inches(7.0), right_heading, right_body),
        ]:
            # Column heading
            self._add_rectangle(slide,
                                 x=x_offset, y=Inches(1.4), w=col_w, h=Inches(0.5),
                                 fill_color=self.palette["secondary"])
            self._add_text_box(slide, heading,
                                x=x_offset + Inches(0.1), y=Inches(1.4),
                                w=col_w - Inches(0.2), h=Inches(0.5),
                                font_size=16, bold=True,
                                color=self.palette["text_dark"])

            # Column body
            body_text = "\n".join(f"•  {i}" for i in body) if isinstance(body, list) else body
            self._add_text_box(slide, body_text,
                                x=x_offset, y=Inches(2.0), w=col_w, h=Inches(4.8),
                                font_size=16, color=self.palette["text_dark"], wrap=True)

        return self

    def stat_slide(self, title: str, stats: list[tuple]) -> "PptxBuilder":
        """
        Big-number stat callout slide.
        stats: list of (number, label) tuples, e.g. [("$2.3B", "AUM"), ("14%", "Vol")]
        Max 4 stats recommended.
        """
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._fill_background(slide, self.palette["primary"])

        self._add_text_box(slide, title,
                           x=Inches(0.8), y=Inches(0.3), w=Inches(11.7), h=Inches(0.9),
                           font_size=28, bold=True, color=self.palette["text_light"])

        n = len(stats)
        box_w = Inches(10.0 / n)
        x_start = Inches((13.33 - 10.0) / 2)

        for i, (number, label) in enumerate(stats):
            x = x_start + i * box_w
            # Number
            self._add_text_box(slide, str(number),
                                x=x, y=Inches(2.2), w=box_w, h=Inches(2.0),
                                font_size=60, bold=True,
                                color=self.palette["secondary"],
                                align=PP_ALIGN.CENTER)
            # Label
            self._add_text_box(slide, label,
                                x=x, y=Inches(4.3), w=box_w, h=Inches(0.6),
                                font_size=18,
                                color=self.palette["text_light"],
                                align=PP_ALIGN.CENTER)
        return self

    def section_divider(self, section_title: str) -> "PptxBuilder":
        """Minimal section break slide."""
        slide = self.prs.slides.add_slide(self._blank_layout)
        self._fill_background(slide, self.palette["secondary"])

        self._add_text_box(slide, section_title,
                           x=Inches(1.5), y=Inches(2.8), w=Inches(10.3), h=Inches(1.8),
                           font_size=40, bold=True,
                           color=self.palette["text_dark"],
                           align=PP_ALIGN.CENTER)
        return self

    def ai_label(self, label: str = "AI GENERATED — verify before use") -> "PptxBuilder":
        """Add a red warning badge to the most recently added slide for AI-synthesised content."""
        slide = self.prs.slides[-1]
        self._add_rectangle(slide,
                             x=Inches(9.2), y=Inches(6.9),
                             w=Inches(3.85), h=Inches(0.38),
                             fill_color="C0392B")
        self._add_text_box(slide, label,
                           x=Inches(9.25), y=Inches(6.92),
                           w=Inches(3.75), h=Inches(0.34),
                           font_size=8, bold=True, color="FFFFFF",
                           align=PP_ALIGN.CENTER)
        return self

    def source_label(self, source: str, x=0.4, width=8.0) -> "PptxBuilder":
        """Add a muted source attribution label to the most recently added slide."""
        slide = self.prs.slides[-1]
        self._add_text_box(slide, f"Source: {source}",
                           x=Inches(x), y=Inches(6.97),
                           w=Inches(width), h=Inches(0.28),
                           font_size=7, italic=True,
                           color=self.palette["muted_text"],
                           align=PP_ALIGN.LEFT)
        return self

    def speaker_notes(self, notes: str) -> "PptxBuilder":
        """Set speaker notes on the most recently added slide."""
        slide = self.prs.slides[-1]
        notes_slide = slide.notes_slide
        tf = notes_slide.notes_text_frame
        tf.text = notes
        return self

    # ── low-level helpers ─────────────────────────────────────────────────────

    def _fill_background(self, slide, hex_color: str):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = hex_to_rgb(hex_color)

    def _add_text_box(self, slide, text: str,
                      x, y, w, h,
                      font_size=18, bold=False, italic=False,
                      color=None, align=PP_ALIGN.LEFT, wrap=True, auto_fit=False):
        txBox = slide.shapes.add_textbox(x, y, w, h)
        tf = txBox.text_frame
        tf.word_wrap = wrap
        if auto_fit:
            tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE

        # Clear default paragraph and set text
        p = tf.paragraphs[0]
        p.alignment = align
        run = p.add_run()
        run.text = text
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = hex_to_rgb(color or self.palette["text_dark"])
        return txBox

    def _add_rectangle(self, slide, x, y, w, h, fill_color: str, line_color=None):
        shape = slide.shapes.add_shape(
            1,  # MSO_SHAPE_TYPE.RECTANGLE
            x, y, w, h
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
        if line_color:
            shape.line.color.rgb = hex_to_rgb(line_color)
        else:
            shape.line.fill.background()  # no border
        return shape

    def save(self, output_path: str, validate=True, final=False, report_path=None) -> Path:
        """Save the presentation and run structural QA unless explicitly disabled."""
        path = Path(output_path)
        path.parent.mkdir(parents=True, exist_ok=True)
        self.prs.save(path)
        if validate:
            from pptx_qa import ensure_pptx_quality
            ensure_pptx_quality(path, final=final, report_path=report_path)
        return path


# ── demo ──────────────────────────────────────────────────────────────────────

def demo():
    b = PptxBuilder(palette="midnight_executive")

    b.title_slide(
        title="Portfolio Risk Review",
        subtitle="Q2 2025 | BlackRock Portfolio Risk Modelling"
    )
    b.stat_slide(
        title="Quarter at a Glance",
        stats=[("12%", "Equity Vol ↑"), ("-18bps", "Credit Spreads"), ("0.94", "Sharpe Ratio"), ("$2.3B", "AUM")]
    )
    b.content_slide(
        title="Key Risk Factors",
        body=[
            "Equity beta: outperformed benchmark by 32bps",
            "Duration: rates moved 25bps higher across the curve",
            "Credit: IG tightened, HY widened — net neutral",
            "FX: carry positions contributed +8bps",
        ]
    )
    b.two_column_slide(
        title="Scenario Analysis",
        left_heading="Bear Case",
        left_body=["Equity -15%", "Spreads +120bps", "Rates -50bps"],
        right_heading="Bull Case",
        right_body=["Equity +12%", "Spreads -40bps", "Rates +25bps"],
    )
    b.section_divider("Appendix")
    b.content_slide(
        title="Methodology Notes",
        body="Full revaluation used for non-linear instruments. "
             "Factor exposures computed daily using MSCI Barra models.",
        dark_bg=True,
    )

    out = b.save("demo_output.pptx")
    print(f"Created: {out}")
    print(f"Slides: {len(b.prs.slides)}")


if __name__ == "__main__":
    demo()
