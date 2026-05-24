"""Lightweight quality checks for PowerPoint files built with python-pptx.

The checks in this module intentionally avoid a renderer so they can run during
every chat-generated deck build. Definite delivery issues are errors; visual
signals that may require judgement are warnings for the assistant to review.
"""

from __future__ import annotations

import argparse
import json
import math
import re
from io import BytesIO
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.enum.dml import MSO_FILL


EMU_PER_INCH = 914400
PANEL_MIN_WIDTH = int(0.85 * EMU_PER_INCH)
PANEL_MIN_HEIGHT = int(0.55 * EMU_PER_INCH)
ALIGNMENT_TOLERANCE = int(0.08 * EMU_PER_INCH)
PANEL_SIZE_TOLERANCE = int(0.08 * EMU_PER_INCH)
PANEL_GAP_TOLERANCE = int(0.12 * EMU_PER_INCH)
PLACEHOLDER_RE = re.compile(
    r"(\[\s*chart\s*\]|lorem ipsum|\bplaceholder\b|\bTODO\b)",
    flags=re.IGNORECASE,
)


def _issue(level: str, slide_no: int | None, check: str, message: str) -> dict:
    issue = {"level": level, "check": check, "message": message}
    if slide_no is not None:
        issue["slide"] = slide_no
    return issue


def _rgb_tuple(value) -> tuple[int, int, int] | None:
    if value is None:
        return None
    text = str(value)
    if len(text) != 6:
        return None
    try:
        return tuple(int(text[i:i + 2], 16) for i in (0, 2, 4))
    except ValueError:
        return None


def _solid_fill_rgb(fill) -> tuple[int, int, int] | None:
    try:
        if fill.type != MSO_FILL.SOLID:
            return None
        return _rgb_tuple(fill.fore_color.rgb)
    except (AttributeError, ValueError):
        return None


def _luminance(rgb: tuple[int, int, int]) -> float:
    channels = []
    for value in rgb:
        channel = value / 255.0
        channels.append(channel / 12.92 if channel <= 0.03928 else ((channel + 0.055) / 1.055) ** 2.4)
    return 0.2126 * channels[0] + 0.7152 * channels[1] + 0.0722 * channels[2]


def _contrast_ratio(foreground: tuple[int, int, int], background: tuple[int, int, int]) -> float:
    light, dark = sorted((_luminance(foreground), _luminance(background)), reverse=True)
    return (light + 0.05) / (dark + 0.05)


def _shape_text(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        return shape.text.strip()
    if getattr(shape, "has_table", False):
        texts = []
        for row in shape.table.rows:
            texts.extend(cell.text.strip() for cell in row.cells if cell.text.strip())
        return " ".join(texts)
    return ""


def _shape_font_sizes(shape) -> list[float]:
    sizes = []
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.font.size is not None:
                    sizes.append(run.font.size.pt)
    return sizes


def _text_colors(shape) -> list[tuple[int, int, int]]:
    colors = []
    if not getattr(shape, "has_text_frame", False):
        return colors
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            try:
                color = _rgb_tuple(run.font.color.rgb)
            except (AttributeError, ValueError):
                color = None
            if color:
                colors.append(color)
    return colors


def _line_estimate(shape, text: str, font_size: float) -> tuple[int, int]:
    width_inches = max(shape.width / EMU_PER_INCH, 0.1)
    height_inches = max(shape.height / EMU_PER_INCH, 0.1)
    chars_per_line = max(int(width_inches * 72 / (font_size * 0.52)), 1)
    estimated_lines = sum(
        max(1, math.ceil(len(line) / chars_per_line))
        for line in (text.splitlines() or [""])
    )
    capacity = max(int(height_inches * 72 / (font_size * 1.25)), 1)
    return estimated_lines, capacity


def _background_beneath_text(shapes, shape_index, shape, slide_background):
    centre_x = shape.left + shape.width / 2
    centre_y = shape.top + shape.height / 2
    text_area = max(shape.width * shape.height, 1)
    for candidate in reversed(shapes[:shape_index]):
        if (
            candidate.left <= centre_x <= candidate.left + candidate.width
            and candidate.top <= centre_y <= candidate.top + candidate.height
        ):
            overlap_width = max(
                0,
                min(shape.left + shape.width, candidate.left + candidate.width)
                - max(shape.left, candidate.left),
            )
            overlap_height = max(
                0,
                min(shape.top + shape.height, candidate.top + candidate.height)
                - max(shape.top, candidate.top),
            )
            if overlap_width * overlap_height / text_area < 0.25:
                continue
            color = _solid_fill_rgb(getattr(candidate, "fill", None))
            if color:
                return color
    return slide_background


def _is_footer_shape(shape, slide_height: int) -> bool:
    """True for shapes in the bottom 9% of the slide (source labels, footnotes, ai labels)."""
    return shape.top > slide_height * 0.91


def _has_visual_element(shapes, slide_height: int) -> bool:
    """True if slide has a table, image, native chart, full-height panel, or >=3 decorative shapes."""
    if any(
        hasattr(s, "image")
        or getattr(s, "has_table", False)
        or getattr(s, "has_chart", False)
        for s in shapes
    ):
        return True
    # Full-height solid panels used in recommendation and assertion_evidence layouts
    for s in shapes:
        if not (getattr(s, "text", "") or "").strip():
            if s.height > slide_height * 0.50:
                try:
                    if s.fill.type == MSO_FILL.SOLID:
                        return True
                except (AttributeError, Exception):
                    pass
    # Pattern slides (status, timeline, process, etc.) use many filled shapes with no text content
    decorative = sum(
        1 for s in shapes
        if not hasattr(s, "image")
        and not getattr(s, "has_table", False)
        and not getattr(s, "has_chart", False)
        and not (getattr(s, "text", "") or "").strip()
    )
    return decorative >= 3


def _is_text_shape(shape) -> bool:
    """True if shape has visible text content (not just a structural/decorative shape)."""
    return getattr(shape, "has_text_frame", False) and bool((getattr(shape, "text", "") or "").strip())


def _shapes_overlap(a, b) -> bool:
    """True if two shapes have overlapping bounding boxes."""
    return (
        a.left < b.left + b.width
        and a.left + a.width > b.left
        and a.top < b.top + b.height
        and a.top + a.height > b.top
    )


def _aligned_clusters(shapes, axis: str) -> list[list]:
    """Group panel-like shapes that share a visual row or column."""
    coordinate = (lambda shape: shape.top) if axis == "row" else (lambda shape: shape.left)
    clusters = []
    for shape in sorted(shapes, key=coordinate):
        position = coordinate(shape)
        for cluster in clusters:
            if abs(position - coordinate(cluster[0])) <= ALIGNMENT_TOLERANCE:
                cluster.append(shape)
                break
        else:
            clusters.append([shape])
    return [cluster for cluster in clusters if len(cluster) >= 3]


def _repeated_panel_issues(shapes, slide_width: int, slide_height: int) -> list[str]:
    """Find uneven rows/columns of repeated large panels without judging bars or decoration."""
    panels = []
    max_area = slide_width * slide_height * 0.45
    for shape in shapes:
        if (
            _shape_text(shape)
            or getattr(shape, "has_table", False)
            or getattr(shape, "has_chart", False)
            or hasattr(shape, "image")
            or shape.width < PANEL_MIN_WIDTH
            or shape.height < PANEL_MIN_HEIGHT
            or shape.width * shape.height > max_area
        ):
            continue
        color = _solid_fill_rgb(getattr(shape, "fill", None))
        if not color:
            continue
        panels.append(shape)

    issues = []
    seen = set()
    if len(panels) < 3:
        return issues
    for axis in ("row", "column"):
        for cluster in _aligned_clusters(panels, axis):
            widths = [shape.width for shape in cluster]
            heights = [shape.height for shape in cluster]
            dimension_problem = (
                max(widths) - min(widths) > PANEL_SIZE_TOLERANCE
                or max(heights) - min(heights) > PANEL_SIZE_TOLERANCE
            )
            ordered = sorted(cluster, key=(lambda shape: shape.left) if axis == "row" else (lambda shape: shape.top))
            gaps = [
                (
                    ordered[index + 1].left - (ordered[index].left + ordered[index].width)
                    if axis == "row"
                    else ordered[index + 1].top - (ordered[index].top + ordered[index].height)
                )
                for index in range(len(ordered) - 1)
            ]
            spacing_problem = gaps and max(gaps) - min(gaps) > PANEL_GAP_TOLERANCE
            if not (dimension_problem or spacing_problem):
                continue
            problem = []
            if dimension_problem:
                problem.append("unequal panel dimensions")
            if spacing_problem:
                problem.append("uneven spacing")
            key = (axis, tuple(sorted(problem)), tuple(id(shape) for shape in ordered))
            if key not in seen:
                seen.add(key)
                issues.append(f"Repeated {axis} panels have {' and '.join(problem)}.")
    return issues


def audit_pptx(path: str | Path, final: bool = False) -> dict:
    """Audit one PPTX file and return a JSON-serializable result dictionary."""
    pptx_path = Path(path)
    errors = []
    warnings = []
    checks = {
        "bounds": 0,
        "placeholder_text": 0,
        "font_size": 0,
        "text_crowding": 0,
        "solid_contrast": 0,
        "image_resolution": 0,
        "table_density": 0,
        "bullet_count": 0,
        "bullet_length": 0,
        "slide_text_volume": 0,
        "visual_density": 0,
        "text_only_font": 0,
        "space_utilisation": 0,
        "shape_overlap": 0,
        "layout_symmetry": 0,
    }

    if not pptx_path.exists():
        errors.append(_issue("error", None, "file", f"File not found: {pptx_path}"))
        return {
            "path": str(pptx_path),
            "slide_count": 0,
            "final": final,
            "passed": False,
            "errors": errors,
            "warnings": warnings,
            "checks": checks,
        }

    try:
        prs = Presentation(str(pptx_path))
    except Exception as exc:
        errors.append(_issue("error", None, "file", f"Could not open PPTX: {exc}"))
        return {
            "path": str(pptx_path),
            "slide_count": 0,
            "final": final,
            "passed": False,
            "errors": errors,
            "warnings": warnings,
            "checks": checks,
        }

    if not prs.slides:
        errors.append(_issue("error", None, "slides", "Presentation has no slides."))

    for slide_no, slide in enumerate(prs.slides, start=1):
        slide_background = _solid_fill_rgb(slide.background.fill)
        shapes = list(slide.shapes)
        slide_word_count = 0
        for shape_index, shape in enumerate(shapes):
            shape_name = getattr(shape, "name", "shape")
            is_footer = _is_footer_shape(shape, prs.slide_height)
            checks["bounds"] += 1
            if (
                shape.left < 0
                or shape.top < 0
                or shape.left + shape.width > prs.slide_width
                or shape.top + shape.height > prs.slide_height
            ):
                errors.append(
                    _issue("error", slide_no, "bounds", f"{shape_name} extends outside the slide canvas.")
                )

            text = _shape_text(shape)
            if text:
                if not is_footer:
                    slide_word_count += len(text.split())
                checks["placeholder_text"] += 1
                if PLACEHOLDER_RE.search(text):
                    item = _issue(
                        "error" if final else "warning",
                        slide_no,
                        "placeholder_text",
                        f"{shape_name} contains visible placeholder text.",
                    )
                    (errors if final else warnings).append(item)

                sizes = _shape_font_sizes(shape)
                if sizes:
                    minimum_size = min(sizes)
                    if not is_footer:
                        checks["font_size"] += 1
                        if minimum_size < 8:
                            warnings.append(
                                _issue(
                                    "warning",
                                    slide_no,
                                    "font_size",
                                    f"{shape_name} uses very small type ({minimum_size:g} pt).",
                                )
                            )
                    estimated_lines, capacity = _line_estimate(shape, text, minimum_size)
                    checks["text_crowding"] += 1
                    if estimated_lines > capacity + 1:
                        severe = final and estimated_lines > capacity * 1.5
                        item = _issue(
                            "error" if severe else "warning",
                            slide_no,
                            "text_crowding",
                            f"{shape_name} text overflow: ~{estimated_lines} lines in space for {capacity}.",
                        )
                        (errors if severe else warnings).append(item)

                text_colors = _text_colors(shape)
                background = _background_beneath_text(shapes, shape_index, shape, slide_background)
                if background and text_colors:
                    checks["solid_contrast"] += 1
                    lowest_ratio = min(_contrast_ratio(color, background) for color in text_colors)
                    largest_font = max(sizes) if sizes else 12
                    threshold = 3.0 if largest_font >= 18 else 4.5
                    if lowest_ratio < threshold:
                        warnings.append(
                            _issue(
                                "warning",
                                slide_no,
                                "solid_contrast",
                                f"{shape_name} has low text contrast ({lowest_ratio:.1f}:1).",
                            )
                        )

                bullet_count = text.count("•")
                if bullet_count > 5:
                    checks["bullet_count"] += 1
                    warnings.append(_issue(
                        "warning", slide_no, "bullet_count",
                        f"{shape_name} has {bullet_count} bullet points - aim for 5 or fewer per slide.",
                    ))
                if bullet_count >= 2:
                    lines = [ln for ln in text.split("•") if ln.strip()]
                    avg_words = sum(len(ln.split()) for ln in lines) / len(lines)
                    if avg_words > 15:
                        checks["bullet_length"] += 1
                        warnings.append(_issue(
                            "warning", slide_no, "bullet_length",
                            f"{shape_name} averages {avg_words:.0f} words per bullet - aim for <=15 words each.",
                        ))

            if getattr(shape, "has_table", False):
                checks["table_density"] += 1
                if len(shape.table.rows) > 12 or len(shape.table.columns) > 8:
                    warnings.append(
                        _issue(
                            "warning",
                            slide_no,
                            "table_density",
                            f"{shape_name} is dense ({len(shape.table.rows)} rows x {len(shape.table.columns)} columns).",
                        )
                    )

            if hasattr(shape, "image"):
                checks["image_resolution"] += 1
                try:
                    with Image.open(BytesIO(shape.image.blob)) as image:
                        dpi_x = image.width / max(shape.width / EMU_PER_INCH, 0.01)
                        dpi_y = image.height / max(shape.height / EMU_PER_INCH, 0.01)
                    if min(dpi_x, dpi_y) < 96:
                        warnings.append(
                            _issue(
                                "warning",
                                slide_no,
                                "image_resolution",
                                f"{shape_name} renders below 96 DPI ({min(dpi_x, dpi_y):.0f} DPI).",
                            )
                        )
                except Exception:
                    warnings.append(
                        _issue("warning", slide_no, "image_resolution", f"Could not inspect {shape_name} resolution.")
                    )

        checks["slide_text_volume"] += 1
        if slide_word_count > 80:
            warnings.append(_issue(
                "warning", slide_no, "slide_text_volume",
                f"Slide {slide_no} has ~{slide_word_count} words - aim for <=80 words total per content slide.",
            ))

        has_body_text = slide_word_count > 15
        if has_body_text:
            has_visual = _has_visual_element(shapes, prs.slide_height)
            if not has_visual:
                checks["visual_density"] += 1
                warnings.append(_issue(
                    "warning", slide_no, "visual_density",
                    f"Slide {slide_no} is text-only - add a table, chart image, or visual pattern.",
                ))
                body_sizes = [sz for s in shapes if not _is_footer_shape(s, prs.slide_height) for sz in _shape_font_sizes(s)]
                if body_sizes and min(body_sizes) < 18:
                    checks["text_only_font"] += 1
                    warnings.append(_issue(
                        "warning", slide_no, "text_only_font",
                        f"Slide {slide_no} is text-only with small type ({min(body_sizes):.0f}pt min) - increase body font to >=18pt.",
                    ))
            checks["space_utilisation"] += 1
            if not has_visual and slide_word_count < 40:
                warnings.append(_issue(
                    "warning", slide_no, "space_utilisation",
                    f"Slide {slide_no} has ~{slide_word_count} words with no visual - canvas will appear empty.",
                ))

        text_shapes = [s for s in shapes if _is_text_shape(s)]
        for idx_a in range(len(text_shapes)):
            for idx_b in range(idx_a + 1, len(text_shapes)):
                sa, sb = text_shapes[idx_a], text_shapes[idx_b]
                if _shapes_overlap(sa, sb):
                    ow = max(0, min(sa.left + sa.width, sb.left + sb.width) - max(sa.left, sb.left))
                    oh = max(0, min(sa.top + sa.height, sb.top + sb.height) - max(sa.top, sb.top))
                    min_area = min(sa.width * sa.height, sb.width * sb.height)
                    if min_area > 0 and ow * oh / min_area > 0.05:
                        checks["shape_overlap"] += 1
                        warnings.append(_issue(
                            "warning", slide_no, "shape_overlap",
                            f"{getattr(sa, 'name', 'shape')} and {getattr(sb, 'name', 'shape')} text boxes overlap - content may render on top of each other.",
                        ))

        for message in _repeated_panel_issues(shapes, prs.slide_width, prs.slide_height):
            checks["layout_symmetry"] += 1
            warnings.append(_issue("warning", slide_no, "layout_symmetry", message))

    return {
        "path": str(pptx_path),
        "slide_count": len(prs.slides),
        "final": final,
        "passed": not errors,
        "errors": errors,
        "warnings": warnings,
        "checks": checks,
    }


def _print_qa_summary(report: dict) -> None:
    """Print a human-readable QA summary to stdout."""
    import sys
    path = report.get("path", "")
    slide_count = report.get("slide_count", 0)
    errors = report.get("errors", [])
    warnings = report.get("warnings", [])

    status = "FAILED" if errors else ("PASSED" if not warnings else "PASSED (warnings)")
    sep = "-" * 60

    def out(text=""):
        try:
            print(text, file=sys.stdout)
        except UnicodeEncodeError:
            print(text.encode("ascii", errors="replace").decode("ascii"), file=sys.stdout)

    out(f"\n{sep}")
    out(f"  QA {status}")
    out(f"  File   : {path}")
    out(f"  Slides : {slide_count}  |  Errors: {len(errors)}  |  Warnings: {len(warnings)}")
    out(sep)
    for issue in errors:
        slide_ref = f"[slide {issue['slide']}] " if "slide" in issue else ""
        out(f"  ERROR   {slide_ref}{issue['message']}")
    for issue in warnings:
        slide_ref = f"[slide {issue['slide']}] " if "slide" in issue else ""
        out(f"  WARN    {slide_ref}{issue['message']}")
    if not errors and not warnings:
        out("  All checks clean.")
    out(f"{sep}\n")


def ensure_pptx_quality(
    path: str | Path,
    final: bool = False,
    report_path: str | Path | None = None,
) -> dict:
    """Audit a PPTX, optionally write the report, and fail on blocking issues."""
    report = audit_pptx(path, final=final)
    if report_path:
        output = Path(report_path)
        output.parent.mkdir(parents=True, exist_ok=True)
        output.write_text(json.dumps(report, indent=2), encoding="utf-8")
    _print_qa_summary(report)
    if report["errors"]:
        messages = "; ".join(issue["message"] for issue in report["errors"])
        raise RuntimeError(f"PPTX quality gate failed: {messages}")
    return report


def main() -> int:
    parser = argparse.ArgumentParser(description="Audit a PPTX before delivery.")
    parser.add_argument("pptx_path", help="PowerPoint file to check")
    parser.add_argument("--final", action="store_true", help="Fail if visible placeholder content remains")
    parser.add_argument("--report", dest="report_path", help="Optional JSON report output path")
    args = parser.parse_args()

    report = audit_pptx(args.pptx_path, final=args.final)
    if args.report_path:
        report_path = Path(args.report_path)
        report_path.parent.mkdir(parents=True, exist_ok=True)
        report_path.write_text(json.dumps(report, indent=2), encoding="utf-8")

    if report["errors"]:
        for issue in report["errors"]:
            print(f"ERROR: slide {issue.get('slide', '-')}: {issue['message']}")
        return 1
    print(
        f"QA passed: {report['slide_count']} slides, "
        f"{len(report['warnings'])} warning(s), {args.pptx_path}"
    )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
