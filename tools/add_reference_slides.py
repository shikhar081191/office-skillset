"""
tools/add_reference_slides.py
Extract reference slides from a PPTX and add them to skills/REFERENCE_GALLERY.md.

Usage:
    python tools/add_reference_slides.py <pptx_path> [slide_numbers]

    slide_numbers: comma-separated, 1-based (e.g. 1,2,3,4,5,6,7,8)
                   omit to extract all slides

Examples:
    python tools/add_reference_slides.py C:/Users/DELL/Downloads/AI_MDS_Project.pptx 1,2,3,4,5,6,7,8
    python tools/add_reference_slides.py decks/reference.pptx
"""

import sys
import os
import re
from pathlib import Path

ROOT = Path(__file__).parent.parent
sys.path.insert(0, str(ROOT))

from pptx import Presentation
from pptx.util import Pt

GALLERY_DIR = ROOT / "skills" / "examples"
GALLERY_MD  = ROOT / "skills" / "REFERENCE_GALLERY.md"


# ── PNG export ─────────────────────────────────────────────────────────────────

def export_pngs(pptx_path: Path, slide_indices: list[int], out_dir: Path) -> dict[int, Path]:
    """Export slides as PNG using PowerPoint COM. Returns {1-based index -> png path}."""
    out_dir.mkdir(parents=True, exist_ok=True)
    try:
        import win32com.client
        import pythoncom
        pythoncom.CoInitialize()
        ppt = win32com.client.Dispatch("PowerPoint.Application")
        ppt.Visible = 1
        abs_path = str(pptx_path.resolve()).replace("/", "\\")
        deck = ppt.Presentations.Open(abs_path, ReadOnly=True, Untitled=False, WithWindow=False)
        results = {}
        for idx in slide_indices:
            slide = deck.Slides(idx)
            out_path = out_dir / f"ref_slide_{idx:02d}.png"
            slide.Export(str(out_path.resolve()), "PNG", 1400, 788)
            results[idx] = out_path
            print(f"  Exported slide {idx} → {out_path.relative_to(ROOT)}")
        deck.Close()
        ppt.Quit()
        return results
    except ImportError:
        print("win32com not available — falling back to manual export instructions.")
        print("In PowerPoint: File → Export → Change File Type → PNG → Save As")
        print(f"Save each slide to: {out_dir}")
        return {}


# ── Metadata extraction ────────────────────────────────────────────────────────

def _rgb_str(shape):
    try:
        rgb = shape.fill.fore_color.rgb
        return f"#{rgb}"
    except Exception:
        return None


def extract_slide_meta(slide, slide_idx: int) -> dict:
    """Extract font sizes, colours, word count, shape summary from a slide."""
    font_sizes = set()
    word_count = 0
    bg_colors  = []
    text_shapes = 0
    fill_shapes = 0

    for s in slide.shapes:
        txt = (getattr(s, "text", "") or "").strip()
        if txt:
            text_shapes += 1
            word_count += len(txt.split())
            if hasattr(s, "text_frame"):
                for para in s.text_frame.paragraphs:
                    for run in para.runs:
                        if run.font.size:
                            font_sizes.add(round(run.font.size.pt))
        else:
            rgb = _rgb_str(s)
            if rgb and rgb != "#None":
                bg_colors.append(rgb)
                fill_shapes += 1

    return {
        "slide":        slide_idx,
        "word_count":   word_count,
        "font_sizes":   sorted(font_sizes, reverse=True),
        "bg_colors":    list(dict.fromkeys(bg_colors)),  # deduplicated, ordered
        "text_shapes":  text_shapes,
        "fill_shapes":  fill_shapes,
        "total_shapes": len(slide.shapes),
    }


# ── Markdown generation ────────────────────────────────────────────────────────

# Describes what makes each slide good — edit as needed per deck
SLIDE_NOTES = {
    1: ("title_slide",
        "Professional title: two-line large title (44 pt), subtitle (16 pt), footer (11 pt). "
        "Decorative teal circle adds visual weight without clutter. Left accent bar anchors layout."),
    2: ("problem_slide / assertion_evidence_slide",
        "28 pt title on dark header bar. Two-column split: left=bullets with coloured dot markers (13 pt), "
        "right=warm panel listing current process steps (12 pt). Bottom callout bar for the punchline. "
        "77 words — dense but not overcrowded because the two-column layout distributes weight."),
    3: ("three_column_card_slide",
        "26 pt title on teal header. Dark quote/input bar spans full width (15 pt). "
        "Three equal white cards below, each with a 4 px teal top rule, 13 pt card title, 12 pt body. "
        "Cards fill the bottom 60 % of the slide — no vacant space."),
    4: ("numbered_steps_slide",
        "26 pt title on dark header. Numbered circles (18 pt, teal/amber) left-aligned. "
        "Step titles 13 pt bold; descriptions 11 pt. Connector lines between steps. "
        "Right sidebar callout box highlights the critical step. Every row is tight — no padding waste."),
    5: ("human_loop / two_column_contrast_slide",
        "Good space coverage: icon-led feature blocks fill both columns. "
        "Multiple visual elements (icons + text) prevent text-only feel. Professional and readable."),
    6: ("status_slide / comparison_slide",
        "Great colour diversity: teal, navy, amber used deliberately. "
        "Clear visual hierarchy with contrasting backgrounds per section."),
    7: ("kpi_dashboard_slide / numbers_slide",
        "Large KPI numbers dominate — good size hierarchy. "
        "Generous white space around numbers makes them breathe. Labels clearly subordinate."),
    8: ("four_box_slide / two_by_two_slide",
        "Title + subtitle header. Four symmetrical boxes in a 2×2 grid. "
        "Consistent padding, aligned edges. Each box has heading + short body — no runaway text."),
}


def build_entry(meta: dict, png_rel: str | None, notes: tuple) -> str:
    pattern_name, description = notes
    img_line = f"![Slide {meta['slide']}]({png_rel})\n\n" if png_rel else ""
    colors_str = "  ".join(meta["bg_colors"][:5]) if meta["bg_colors"] else "—"
    fonts_str  = ", ".join(f"{f}pt" for f in meta["font_sizes"]) if meta["font_sizes"] else "inherited"

    return (
        f"### Slide {meta['slide']} — {pattern_name}\n\n"
        f"{img_line}"
        f"**What makes it work:** {description}\n\n"
        f"| Metric | Value |\n"
        f"|--------|-------|\n"
        f"| Word count | {meta['word_count']} |\n"
        f"| Font sizes used | {fonts_str} |\n"
        f"| Fill colours | `{colors_str}` |\n"
        f"| Shapes total | {meta['total_shapes']} ({meta['text_shapes']} text, {meta['fill_shapes']} filled) |\n\n"
    )


def update_gallery(entries: list[str]) -> None:
    GALLERY_DIR.mkdir(parents=True, exist_ok=True)

    header = (
        "# Reference Gallery — Good Slide Examples\n\n"
        "These slides are extracted from reference decks approved as high-quality.\n"
        "Read this file before choosing patterns or setting font sizes.\n"
        "Key lessons drawn from each slide are in the 'What makes it work' field.\n\n"
        "---\n\n"
    )

    existing = ""
    if GALLERY_MD.exists():
        existing_text = GALLERY_MD.read_text(encoding="utf-8")
        # Keep everything after the header separator
        split = existing_text.split("---\n\n", 1)
        existing = split[1] if len(split) > 1 else ""

    # Merge: replace existing entries for same slide number, append new ones
    new_nums = {int(re.search(r"### Slide (\d+)", e).group(1)) for e in entries if re.search(r"### Slide (\d+)", e)}
    kept = []
    for block in re.split(r"(?=### Slide \d+)", existing):
        m = re.search(r"### Slide (\d+)", block)
        if m and int(m.group(1)) in new_nums:
            continue  # will be replaced
        if block.strip():
            kept.append(block)

    all_entries = kept + entries
    # Sort by slide number
    def _num(e):
        m = re.search(r"### Slide (\d+)", e)
        return int(m.group(1)) if m else 999
    all_entries.sort(key=_num)

    GALLERY_MD.write_text(header + "".join(all_entries), encoding="utf-8")
    print(f"\nGallery written: {GALLERY_MD.relative_to(ROOT)}")


# ── AI_INSTRUCTIONS hook ───────────────────────────────────────────────────────

def ensure_instructions_reference() -> None:
    """Add a reference to REFERENCE_GALLERY.md in AI_INSTRUCTIONS.md if not present."""
    ai_md = ROOT / "AI_INSTRUCTIONS.md"
    if not ai_md.exists():
        return
    text = ai_md.read_text(encoding="utf-8")
    if "REFERENCE_GALLERY" in text:
        return
    injection = (
        "3. `skills/REFERENCE_GALLERY.md` — visual examples of good slides with annotations\n"
        "   explaining what works: font hierarchy, space coverage, colour use.\n"
        "   Match the style of these examples when choosing layouts and font sizes.\n\n"
    )
    # Insert after the existing item 2 in Step 0
    text = text.replace(
        "2. `skills/` folder — individual pattern files for any pattern you plan to use\n",
        "2. `skills/` folder — individual pattern files for any pattern you plan to use\n" + injection,
    )
    ai_md.write_text(text, encoding="utf-8")
    print("AI_INSTRUCTIONS.md updated to reference REFERENCE_GALLERY.md")


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    if len(sys.argv) < 2:
        print(__doc__)
        sys.exit(1)

    pptx_path = Path(sys.argv[1])
    if not pptx_path.exists():
        print(f"File not found: {pptx_path}")
        sys.exit(1)

    prs = Presentation(str(pptx_path))
    total = len(prs.slides)

    if len(sys.argv) >= 3:
        slide_indices = [int(x.strip()) for x in sys.argv[2].split(",")]
    else:
        slide_indices = list(range(1, total + 1))

    print(f"Processing {len(slide_indices)} slides from {pptx_path.name} ({total} total)")
    print()

    # Export PNGs
    png_map = export_pngs(pptx_path, slide_indices, GALLERY_DIR)

    # Extract metadata and build gallery entries
    entries = []
    for idx in slide_indices:
        slide = prs.slides[idx - 1]
        meta  = extract_slide_meta(slide, idx)
        png_rel = str(png_map[idx].relative_to(ROOT)).replace("\\", "/") if idx in png_map else None
        notes   = SLIDE_NOTES.get(idx, (f"slide_{idx}", "No annotation provided."))
        entry   = build_entry(meta, png_rel, notes)
        entries.append(entry)
        print(f"  Annotated slide {idx}: {meta['word_count']} words, fonts {meta['font_sizes']}")

    update_gallery(entries)
    ensure_instructions_reference()
    print("\nDone. Run: git add skills/ AI_INSTRUCTIONS.md && git commit -m 'Add reference gallery'")


if __name__ == "__main__":
    main()
