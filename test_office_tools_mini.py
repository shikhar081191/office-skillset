"""
test_office_tools_mini.py — Test suite for office_scripts_mini (Python-only)

No LibreOffice, no npm, no external tools required.
All tests use python-docx, python-pptx, and openpyxl.

Run:
    python test_office_tools_mini.py
    python test_office_tools_mini.py -v
    python test_office_tools_mini.py TestCreateDocx
"""

import json
import shutil
import subprocess
import sys
import tempfile
import unittest
import zipfile
from pathlib import Path


# ─── helpers ──────────────────────────────────────────────────────────────────

def run(args: list, **kwargs) -> subprocess.CompletedProcess:
    return subprocess.run([sys.executable] + args, capture_output=True, text=True, **kwargs)


# ─── test classes ─────────────────────────────────────────────────────────────

class TestImports(unittest.TestCase):
    """All required libraries must be importable."""

    def test_python_docx(self):
        import docx; self.assertTrue(True)

    def test_python_pptx(self):
        import pptx; self.assertTrue(True)

    def test_openpyxl(self):
        import openpyxl; self.assertTrue(True)

    def test_pandas(self):
        import pandas; self.assertTrue(True)

    def test_pillow(self):
        from PIL import Image; self.assertTrue(True)

    def test_defusedxml(self):
        import defusedxml.minidom; self.assertTrue(True)


class TestCreateDocx(unittest.TestCase):
    """Tests for create_docx.py — DocxBuilder"""

    def setUp(self):
        self.tmp = tempfile.mkdtemp()
        self.out = Path(self.tmp) / "output.docx"

    def tearDown(self):
        shutil.rmtree(self.tmp)

    def _builder(self):
        from create_docx import DocxBuilder
        return DocxBuilder()

    def test_creates_file(self):
        b = self._builder()
        b.heading("Test", level=1).paragraph("Hello.").save(str(self.out))
        self.assertTrue(self.out.exists())

    def test_output_is_valid_docx(self):
        b = self._builder()
        b.heading("Test").save(str(self.out))
        with zipfile.ZipFile(self.out) as zf:
            self.assertIn("[Content_Types].xml", zf.namelist())

    def test_heading_in_content(self):
        b = self._builder()
        b.heading("My Heading", level=1).save(str(self.out))
        from docx import Document
        doc = Document(self.out)
        texts = [p.text for p in doc.paragraphs]
        self.assertIn("My Heading", texts)

    def test_paragraph_in_content(self):
        b = self._builder()
        b.paragraph("Hello World").save(str(self.out))
        from docx import Document
        doc = Document(self.out)
        texts = " ".join(p.text for p in doc.paragraphs)
        self.assertIn("Hello World", texts)

    def test_bullet_list(self):
        b = self._builder()
        b.bullet_list(["Item A", "Item B", "Item C"]).save(str(self.out))
        from docx import Document
        doc = Document(self.out)
        texts = [p.text for p in doc.paragraphs]
        self.assertIn("Item A", texts)
        self.assertIn("Item C", texts)

    def test_numbered_list(self):
        b = self._builder()
        b.numbered_list(["First", "Second"]).save(str(self.out))
        from docx import Document
        doc = Document(self.out)
        texts = [p.text for p in doc.paragraphs]
        self.assertIn("First", texts)

    def test_table(self):
        b = self._builder()
        b.table([["Header", "Value"], ["Row1", "10"], ["Row2", "20"]]).save(str(self.out))
        from docx import Document
        doc = Document(self.out)
        self.assertEqual(len(doc.tables), 1)
        self.assertEqual(doc.tables[0].cell(0, 0).text, "Header")

    def test_page_size_us_letter(self):
        """Page should be US Letter (8.5 x 11 inches), not A4."""
        b = self._builder()
        b.save(str(self.out))
        from docx import Document
        from docx.shared import Inches
        doc = Document(self.out)
        section = doc.sections[0]
        self.assertAlmostEqual(section.page_width.inches, 8.5, places=1)
        self.assertAlmostEqual(section.page_height.inches, 11.0, places=1)

    def test_demo_runs(self):
        r = run(["create_docx.py"])
        self.assertEqual(r.returncode, 0, f"stderr: {r.stderr}")
        self.assertTrue(Path("demo_output.docx").exists())

    def test_chaining(self):
        """Builder should support method chaining."""
        b = self._builder()
        result = b.heading("H").paragraph("P").bullet_list(["A"]).save(str(self.out))
        self.assertTrue(self.out.exists())

    def test_final_save_writes_quality_report(self):
        report_path = Path(self.tmp) / "output.qa.json"
        b = self._builder()
        b.heading("Checked Note").paragraph("Final document content.")
        b.save(str(self.out), final=True, report_path=report_path)
        report = json.loads(report_path.read_text(encoding="utf-8"))
        self.assertTrue(report["passed"])

    def test_final_save_fails_on_visible_placeholder(self):
        b = self._builder()
        b.heading("Draft").paragraph("TODO: insert conclusion")
        with self.assertRaises(RuntimeError):
            b.save(str(self.out), final=True)


class TestCreatePptx(unittest.TestCase):
    """Tests for create_pptx.py — PptxBuilder"""

    def setUp(self):
        self.tmp = tempfile.mkdtemp()
        self.out = Path(self.tmp) / "output.pptx"

    def tearDown(self):
        shutil.rmtree(self.tmp)

    def _builder(self, palette="midnight_executive"):
        from create_pptx import PptxBuilder
        return PptxBuilder(palette=palette)

    def test_creates_file(self):
        b = self._builder()
        b.title_slide("Test").save(str(self.out))
        self.assertTrue(self.out.exists())

    def test_output_is_valid_pptx(self):
        b = self._builder()
        b.title_slide("Test").save(str(self.out))
        with zipfile.ZipFile(self.out) as zf:
            self.assertTrue(any("presentation.xml" in n for n in zf.namelist()))

    def test_title_slide_adds_one_slide(self):
        b = self._builder()
        b.title_slide("My Title", "My Subtitle")
        self.assertEqual(len(b.prs.slides), 1)

    def test_content_slide(self):
        b = self._builder()
        b.title_slide("Title").content_slide("Content", body=["Point 1", "Point 2"])
        self.assertEqual(len(b.prs.slides), 2)

    def test_two_column_slide(self):
        b = self._builder()
        b.two_column_slide("Title", "Left", ["A", "B"], "Right", ["C", "D"])
        self.assertEqual(len(b.prs.slides), 1)

    def test_stat_slide(self):
        b = self._builder()
        b.stat_slide("Stats", [("$1B", "AUM"), ("12%", "Vol")])
        self.assertEqual(len(b.prs.slides), 1)

    def test_section_divider(self):
        b = self._builder()
        b.section_divider("Section 1")
        self.assertEqual(len(b.prs.slides), 1)

    def test_all_palettes(self):
        from create_pptx import PALETTES
        for palette_name in PALETTES:
            b = self._builder(palette=palette_name)
            out = Path(self.tmp) / f"test_{palette_name}.pptx"
            b.title_slide("Test").save(str(out))
            self.assertTrue(out.exists(), f"Failed for palette: {palette_name}")

    def test_slide_dimensions_widescreen(self):
        """Slides should be 13.33 x 7.5 inches (16:9)."""
        from create_pptx import PptxBuilder
        from pptx.util import Inches
        b = PptxBuilder()
        self.assertAlmostEqual(b.prs.slide_width / Inches(1), 13.33, places=1)
        self.assertAlmostEqual(b.prs.slide_height / Inches(1), 7.5, places=1)

    def test_multi_slide_deck(self):
        b = self._builder()
        b.title_slide("Title")
        b.content_slide("Slide 2", "Body text")
        b.stat_slide("Stats", [("42", "Answer")])
        b.section_divider("End")
        b.save(str(self.out))
        from pptx import Presentation
        prs = Presentation(self.out)
        self.assertEqual(len(prs.slides), 4)

    def test_final_save_writes_quality_report(self):
        b = self._builder()
        report_path = Path(self.tmp) / "output.qa.json"
        b.title_slide("Checked Deck", "Delivery-ready result")
        b.save(str(self.out), final=True, report_path=report_path)
        report = json.loads(report_path.read_text(encoding="utf-8"))
        self.assertTrue(report["passed"])
        self.assertEqual(report["slide_count"], 1)

    def test_save_fails_on_out_of_bounds_shape(self):
        from pptx.util import Inches
        b = self._builder()
        b.title_slide("Invalid Deck")
        slide = b.prs.slides[0]
        b._add_rectangle(
            slide,
            x=Inches(13.0), y=Inches(1.0), w=Inches(1.0), h=Inches(0.3),
            fill_color=b.palette["accent"],
        )
        with self.assertRaises(RuntimeError):
            b.save(str(self.out))

    def test_final_save_fails_on_visible_placeholder(self):
        b = self._builder()
        b.content_slide("Chart", "[ chart placeholder ]")
        with self.assertRaises(RuntimeError):
            b.save(str(self.out), final=True)

    def test_qa_flags_uneven_repeated_panels(self):
        from pptx.util import Inches
        from pptx_qa import audit_pptx
        b = self._builder()
        b.title_slide("Layout Review")
        slide = b.prs.slides[0]
        for x, width in [(1.0, 1.5), (3.0, 1.5), (5.8, 2.0)]:
            b._add_rectangle(
                slide,
                x=Inches(x), y=Inches(2.0), w=Inches(width), h=Inches(0.8),
                fill_color=b.palette["surface"],
            )
        b.save(str(self.out), validate=False)
        report = audit_pptx(self.out, final=True)
        self.assertGreater(report["checks"]["layout_symmetry"], 0)
        self.assertTrue(any(item["check"] == "layout_symmetry" for item in report["warnings"]))

    def test_demo_runs(self):
        r = run(["create_pptx.py"])
        self.assertEqual(r.returncode, 0, f"stderr: {r.stderr}")
        self.assertTrue(Path("demo_output.pptx").exists())


class TestPptxPatterns(unittest.TestCase):
    """Tests for reusable presentation patterns and their validation rules."""

    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp())
        from PIL import Image
        self.chart_path = self.tmp / "chart.png"
        Image.new("RGB", (640, 360), "#FFFFFF").save(self.chart_path)

    def tearDown(self):
        shutil.rmtree(self.tmp)

    def test_existing_builder_palette_is_used_by_patterns(self):
        from create_pptx import PptxBuilder, hex_to_rgb
        from patterns import status_slide
        b = PptxBuilder(palette="teal_trust")
        status_slide("Status", [{"name": "Model", "status": "Ready", "rag": "green"}],
                     builder=b)
        accent_rule = b.prs.slides[0].shapes[1]
        self.assertEqual(accent_rule.fill.fore_color.rgb, hex_to_rgb(b.palette["secondary"]))

    def test_new_patterns_generate_one_slide(self):
        from patterns import (
            annotation_chart_slide, assertion_evidence_slide, diverging_bar_slide,
            heat_map_slide, scorecard_slide, two_by_two_slide, waterfall_slide,
        )
        builders = [
            heat_map_slide("Heat", ["Rates"], ["A", "B"], [[0.3, 0.6]]),
            waterfall_slide("Waterfall", [{"name": "Rates", "value": 3.4}]),
            scorecard_slide("Score", ["Accuracy"], ["A", "B"], [[3, 5]]),
            annotation_chart_slide("Chart", self.chart_path,
                                   [{"label": "Spike", "x": 0.5, "y": 0.4}]),
            two_by_two_slide("Matrix", "Effort", "Impact", "Low", "High", "Low", "High",
                             [{"name": "Alerts", "x": 0.2, "y": 0.8}]),
            assertion_evidence_slide("Finding", "Model B reduces tail error",
                                     [{"text": "PMSE lower across cohorts", "stat": "18%"}]),
            diverging_bar_slide("Bars", ["Rates"], "Current", "Proposed", [0.6], [0.4]),
        ]
        self.assertTrue(all(len(builder.prs.slides) == 1 for builder in builders))

    def test_pattern_validation_rules(self):
        from patterns import annotation_chart_slide, heat_map_slide, scorecard_slide
        with self.assertRaises(ValueError):
            heat_map_slide("Heat", ["Rates"], ["A", "B"], [[0.3]])
        with self.assertRaises(ValueError):
            scorecard_slide("Score", ["Accuracy", "Runtime"], ["A"],
                            [[4], [3]], weights=[0.5, 0.4])
        with self.assertRaises(ValueError):
            annotation_chart_slide("Chart", self.chart_path,
                                   [{"label": "Outside", "x": 1.2, "y": 0.4}])

    def test_pattern_demos_generate_requested_outputs(self):
        from patterns import demo_all
        demo_all()
        for name in [
            "heat_map_slide", "waterfall_slide", "scorecard_slide",
            "annotation_chart_slide", "two_by_two_slide",
            "assertion_evidence_slide", "diverging_bar_slide",
        ]:
            self.assertTrue(Path(f"output/demo_{name}.pptx").exists(), name)


class TestCreateXlsx(unittest.TestCase):
    """Tests for create_xlsx.py — XlsxBuilder"""

    def setUp(self):
        self.tmp = tempfile.mkdtemp()
        self.out = Path(self.tmp) / "output.xlsx"

    def tearDown(self):
        shutil.rmtree(self.tmp)

    def _builder(self):
        from create_xlsx import XlsxBuilder
        return XlsxBuilder()

    def test_creates_file(self):
        b = self._builder()
        b.add_sheet("Sheet1").label_cell(1, 1, "Hello")
        b.save(str(self.out))
        self.assertTrue(self.out.exists())

    def test_output_is_valid_xlsx(self):
        b = self._builder()
        b.add_sheet("Data").label_cell(1, 1, "Test")
        b.save(str(self.out))
        with zipfile.ZipFile(self.out) as zf:
            self.assertTrue(any("workbook.xml" in n for n in zf.namelist()))

    def test_sheet_names(self):
        b = self._builder()
        b.add_sheet("Assumptions")
        b.add_sheet("P&L")
        b.add_sheet("Output")
        b.save(str(self.out))
        from openpyxl import load_workbook
        wb = load_workbook(self.out)
        self.assertIn("Assumptions", wb.sheetnames)
        self.assertIn("P&L", wb.sheetnames)

    def test_input_cell_value(self):
        b = self._builder()
        s = b.add_sheet("Sheet1")
        s.input_cell(1, 1, 42)
        b.save(str(self.out))
        from openpyxl import load_workbook
        wb = load_workbook(self.out, data_only=True)
        ws = wb.active
        self.assertEqual(ws.cell(1, 1).value, 42)

    def test_formula_cell_stored(self):
        b = self._builder()
        s = b.add_sheet("Sheet1")
        s.input_cell(1, 1, 10)
        s.input_cell(2, 1, 20)
        s.formula_cell(3, 1, "=SUM(A1:A2)")
        b.save(str(self.out))
        from openpyxl import load_workbook
        wb = load_workbook(self.out)  # data_only=False to read formula string
        ws = wb.active
        self.assertEqual(ws.cell(3, 1).value, "=SUM(A1:A2)")

    def test_input_cell_blue_font(self):
        """Input cells should have blue font."""
        from create_xlsx import COLOR_INPUT
        b = self._builder()
        s = b.add_sheet("Sheet1")
        s.input_cell(1, 1, 100)
        b.save(str(self.out))
        from openpyxl import load_workbook
        wb = load_workbook(self.out)
        cell = wb.active.cell(1, 1)
        self.assertEqual(cell.font.color.rgb, COLOR_INPUT)

    def test_assumption_cell_yellow_fill(self):
        """Assumption cells should have yellow background."""
        from create_xlsx import COLOR_ASSUMPTION
        b = self._builder()
        s = b.add_sheet("Sheet1")
        s.assumption_cell(1, 1, 0.08)
        b.save(str(self.out))
        from openpyxl import load_workbook
        wb = load_workbook(self.out)
        cell = wb.active.cell(1, 1)
        self.assertEqual(cell.fill.fgColor.rgb, COLOR_ASSUMPTION)

    def test_headers_row(self):
        b = self._builder()
        s = b.add_sheet("Data")
        s.headers(1, ["Name", "Value", "Formula"])
        b.save(str(self.out))
        from openpyxl import load_workbook
        wb = load_workbook(self.out)
        ws = wb.active
        self.assertEqual(ws.cell(1, 1).value, "Name")
        self.assertEqual(ws.cell(1, 3).value, "Formula")

    def test_data_table(self):
        b = self._builder()
        s = b.add_sheet("Sheet1")
        s.data_table(
            start_row=1, start_col=1,
            headers=["A", "B"],
            rows=[["x", 1], ["y", 2]],
        )
        b.save(str(self.out))
        from openpyxl import load_workbook
        wb = load_workbook(self.out)
        ws = wb.active
        self.assertEqual(ws.cell(1, 1).value, "A")
        self.assertEqual(ws.cell(2, 1).value, "x")
        self.assertEqual(ws.cell(3, 2).value, 2)

    def test_multi_sheet_workbook(self):
        b = self._builder()
        b.add_sheet("Sheet1").input_cell(1, 1, "hello")
        b.add_sheet("Sheet2").input_cell(1, 1, "world")
        b.save(str(self.out))
        from openpyxl import load_workbook
        wb = load_workbook(self.out)
        self.assertEqual(len(wb.sheetnames), 2)

    def test_demo_runs(self):
        r = run(["create_xlsx.py"])
        self.assertEqual(r.returncode, 0, f"stderr: {r.stderr}")
        self.assertTrue(Path("demo_output.xlsx").exists())

    def test_final_save_writes_quality_report(self):
        report_path = Path(self.tmp) / "output.qa.json"
        b = self._builder()
        b.add_sheet("Results").label_cell(1, 1, "Complete")
        b.save(str(self.out), final=True, report_path=report_path)
        report = json.loads(report_path.read_text(encoding="utf-8"))
        self.assertTrue(report["passed"])

    def test_final_save_fails_on_invalid_formula_reference(self):
        b = self._builder()
        b.add_sheet("Results").formula_cell(1, 1, "=#REF!+1")
        with self.assertRaises(RuntimeError):
            b.save(str(self.out), final=True)


class TestUnpackPack(unittest.TestCase):
    """Round-trip tests for unpack.py and pack.py"""

    def setUp(self):
        self.tmp = tempfile.mkdtemp()

    def tearDown(self):
        shutil.rmtree(self.tmp)

    def _make_docx(self, name="test.docx"):
        from docx import Document
        path = Path(self.tmp) / name
        doc = Document()
        doc.add_heading("Test", 1)
        doc.add_paragraph("Hello world.")
        doc.save(path)
        return path

    def _make_pptx(self, name="test.pptx"):
        from pptx import Presentation
        path = Path(self.tmp) / name
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Test Slide"
        prs.save(path)
        return path

    def _make_xlsx(self, name="test.xlsx"):
        from openpyxl import Workbook
        path = Path(self.tmp) / name
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Test"
        wb.save(path)
        return path

    def _round_trip(self, src: Path, suffix: str) -> Path:
        out_dir = Path(self.tmp) / "unpacked"
        out_file = Path(self.tmp) / f"repacked.{suffix}"
        r1 = run(["unpack.py", str(src), str(out_dir)])
        self.assertEqual(r1.returncode, 0, f"Unpack: {r1.stdout}")
        r2 = run(["pack.py", str(out_dir), str(out_file)])
        self.assertEqual(r2.returncode, 0, f"Pack: {r2.stdout}")
        return out_file

    def test_roundtrip_docx(self):
        src = self._make_docx()
        out = self._round_trip(src, "docx")
        self.assertTrue(out.exists())
        from docx import Document
        doc = Document(out)
        texts = [p.text for p in doc.paragraphs]
        self.assertIn("Test", texts)

    def test_roundtrip_pptx(self):
        src = self._make_pptx()
        out = self._round_trip(src, "pptx")
        self.assertTrue(out.exists())
        with zipfile.ZipFile(out) as zf:
            self.assertTrue(any("presentation.xml" in n for n in zf.namelist()))

    def test_roundtrip_xlsx(self):
        src = self._make_xlsx()
        out = self._round_trip(src, "xlsx")
        self.assertTrue(out.exists())
        with zipfile.ZipFile(out) as zf:
            self.assertTrue(any("workbook.xml" in n for n in zf.namelist()))

    def test_unpack_pretty_prints_xml(self):
        src = self._make_docx()
        out_dir = Path(self.tmp) / "unpacked"
        run(["unpack.py", str(src), str(out_dir)])
        content = (out_dir / "word" / "document.xml").read_text(encoding="utf-8")
        self.assertIn("\n", content)

    def test_unpack_missing_file_fails(self):
        r = run(["unpack.py", "nonexistent.docx", str(self.tmp)])
        self.assertNotEqual(r.returncode, 0)

    def test_pack_missing_dir_fails(self):
        r = run(["pack.py", "/nonexistent/dir", "output.docx"])
        self.assertNotEqual(r.returncode, 0)


class TestExtractText(unittest.TestCase):
    """Tests for extract_text.py"""

    def setUp(self):
        self.tmp = tempfile.mkdtemp()

    def tearDown(self):
        shutil.rmtree(self.tmp)

    def test_extract_docx(self):
        from docx import Document
        path = Path(self.tmp) / "test.docx"
        doc = Document()
        doc.add_heading("My Heading", 1)
        doc.add_paragraph("My paragraph text.")
        doc.save(path)
        r = run(["extract_text.py", str(path)])
        self.assertEqual(r.returncode, 0)
        self.assertIn("My Heading", r.stdout)
        self.assertIn("My paragraph text", r.stdout)

    def test_extract_pptx(self):
        from pptx import Presentation
        path = Path(self.tmp) / "test.pptx"
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "Slide Title Here"
        prs.save(path)
        r = run(["extract_text.py", str(path)])
        self.assertEqual(r.returncode, 0)
        self.assertIn("Slide Title Here", r.stdout)

    def test_extract_xlsx(self):
        from openpyxl import Workbook
        path = Path(self.tmp) / "test.xlsx"
        wb = Workbook()
        ws = wb.active
        ws["A1"] = "Factor"
        ws["B1"] = "Return"
        ws["A2"] = "Equity"
        wb.save(path)
        r = run(["extract_text.py", str(path)])
        self.assertEqual(r.returncode, 0)
        self.assertIn("Factor", r.stdout)
        self.assertIn("Equity", r.stdout)

    def test_missing_file_fails(self):
        r = run(["extract_text.py", "nonexistent.docx"])
        self.assertNotEqual(r.returncode, 0)


class TestSourceDocx(unittest.TestCase):
    """Tests for structured Word source intake and reusable Office patterns."""

    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp())

    def tearDown(self):
        shutil.rmtree(self.tmp)

    def test_reads_talking_points_and_numeric_table(self):
        from docx import Document
        from source_docx import read_source_docx, source_to_markdown
        path = self.tmp / "source.docx"
        doc = Document()
        doc.add_heading("Treasury Model Results", 1)
        doc.add_paragraph("Bills drive the improvement.")
        table = doc.add_table(rows=3, cols=3)
        for row, values in enumerate([
            ["Cohort", "Current", "Model B"],
            ["Bills", "60.2%", "14.1%"],
            ["Bonds", "3.3%", "0.4%"],
        ]):
            for col, value in enumerate(values):
                table.cell(row, col).text = value
        doc.save(path)
        source = read_source_docx(path)
        self.assertEqual(source["title"], "Treasury Model Results")
        self.assertEqual(source["tables"][0]["numeric_columns"], ["Current", "Model B"])
        self.assertIn("Candidate slide visuals", source_to_markdown(source))

    def test_source_markdown_preserves_table_position_within_section(self):
        from docx import Document
        from source_docx import read_source_docx, source_to_markdown
        path = self.tmp / "ordered_source.docx"
        doc = Document()
        doc.add_heading("Ordered Brief", 1)
        doc.add_paragraph("Narrative before the table.")
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Metric"
        table.cell(0, 1).text = "Value"
        table.cell(1, 0).text = "PMSE"
        table.cell(1, 1).text = "14.1%"
        doc.add_paragraph("Narrative after the table.")
        doc.save(path)

        source = read_source_docx(path)
        markdown = source_to_markdown(source)
        before = markdown.index("Narrative before the table.")
        table_pos = markdown.index("### table_1")
        after = markdown.index("Narrative after the table.")
        self.assertLess(before, table_pos)
        self.assertLess(table_pos, after)
        self.assertEqual(source["tables"][0]["section_heading"], "Ordered Brief")

    def test_source_markdown_captures_header_and_footer_text(self):
        from docx import Document
        from source_docx import read_source_docx, source_to_markdown
        path = self.tmp / "header_footer_source.docx"
        doc = Document()
        doc.add_heading("Brief", 1)
        doc.add_paragraph("Body content.")
        doc.sections[0].header.paragraphs[0].text = "CONFIDENTIAL HEADER MARKER"
        doc.sections[0].footer.paragraphs[0].text = "SOURCE FOOTER MARKER"
        doc.save(path)

        source = read_source_docx(path)
        markdown = source_to_markdown(source)
        self.assertIn("CONFIDENTIAL HEADER MARKER", markdown)
        self.assertIn("SOURCE FOOTER MARKER", markdown)

    def test_research_note_pattern_saves_with_qa(self):
        from docx_patterns import research_note_document
        output = self.tmp / "note.docx"
        report = self.tmp / "note.qa.json"
        b = research_note_document(
            "PMSE Review",
            executive_summary="Model B improves model fit.",
            findings=["Bills show the largest gain."],
        )
        b.save(output, final=True, report_path=report)
        self.assertTrue(output.exists())
        self.assertTrue(json.loads(report.read_text(encoding="utf-8"))["passed"])

    def test_word_source_builds_checked_starter_deck(self):
        from docx import Document
        from build_deck_from_docx import build_deck_from_docx
        path = self.tmp / "brief.docx"
        output = self.tmp / "brief_deck.pptx"
        doc = Document()
        doc.add_heading("Model Review", 1)
        doc.add_paragraph("Model B is recommended for bill risk estimation.")
        table = doc.add_table(rows=3, cols=3)
        for row, values in enumerate([
            ["Cohort", "Current", "Model B"],
            ["Bills", "60.2%", "14.1%"],
            ["Bonds", "3.3%", "0.4%"],
        ]):
            for col, value in enumerate(values):
                table.cell(row, col).text = value
        doc.save(path)
        build_deck_from_docx(path, output)
        self.assertTrue(output.exists())
        report = json.loads(output.with_suffix(".qa.json").read_text(encoding="utf-8"))
        self.assertTrue(report["passed"])

    def test_starter_deck_paginates_without_trimming_text_or_table_cells(self):
        from docx import Document
        from pptx import Presentation
        from build_deck_from_docx import build_deck_from_docx
        path = self.tmp / "complete_brief.docx"
        output = self.tmp / "complete_brief_deck.pptx"
        doc = Document()
        doc.add_heading("Source Fidelity", 1)
        for index in range(8):
            tail = " FINAL_TEXT_MARKER" if index == 7 else ""
            doc.add_paragraph(f"Paragraph {index} " + ("complete context " * 12) + tail)
        table = doc.add_table(rows=11, cols=10)
        for col in range(10):
            table.cell(0, col).text = f"Header_{col}"
        for row in range(1, 11):
            for col in range(10):
                table.cell(row, col).text = f"R{row}C{col}"
        table.cell(10, 9).text = "FINAL_TABLE_MARKER"
        doc.save(path)

        build_deck_from_docx(path, output)
        prs = Presentation(output)
        visible_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if getattr(shape, "has_text_frame", False):
                    visible_text.append(shape.text)
                if getattr(shape, "has_table", False):
                    for row in shape.table.rows:
                        visible_text.extend(cell.text for cell in row.cells)
        rendered_source = "\n".join(visible_text)
        self.assertIn("FINAL_TEXT_MARKER", rendered_source)
        self.assertIn("FINAL_TABLE_MARKER", rendered_source)

    def test_word_source_builds_inside_project_workspace(self):
        from docx import Document
        from build_deck_from_docx import build_deck_from_docx
        project_root = self.tmp / "projects"
        path = self.tmp / "brief.docx"
        doc = Document()
        doc.add_heading("Project Brief", 1)
        doc.add_paragraph("Prepare a model recommendation.")
        doc.save(path)
        output = build_deck_from_docx(
            path,
            project_name="Treasury Review",
            projects_root=project_root,
        )
        project_dir = project_root / "treasury_review"
        self.assertEqual(output, project_dir / "outputs" / "brief_deck.pptx")
        self.assertTrue((project_dir / "inputs" / "brief.docx").exists())
        self.assertTrue((project_dir / "qa" / "brief_deck.qa.json").exists())
        self.assertTrue((project_dir / "project.json").exists())
        self.assertTrue((project_dir / "working" / "brief_source.json").exists())
        self.assertTrue((project_dir / "working" / "brief_source.md").exists())

    def test_unified_workflow_builds_pptx_and_records_manifest(self):
        from docx import Document
        from artifact_workflow import run_workflow
        project_root = self.tmp / "projects"
        path = self.tmp / "model_brief.docx"
        doc = Document()
        doc.add_heading("Treasury Model Recommendation", 1)
        doc.add_paragraph("Model B materially reduces bills PMSE.")
        table = doc.add_table(rows=3, cols=3)
        for row, values in enumerate([
            ["Cohort", "Current", "Model B"],
            ["Bills", "60.2%", "14.1%"],
            ["Bonds", "3.3%", "0.4%"],
        ]):
            for col, value in enumerate(values):
                table.cell(row, col).text = value
        doc.save(path)
        result = run_workflow(
            path,
            project_name="Treasury Decision",
            projects_root=project_root,
            instructions="Focus on the bill improvement.",
        )
        self.assertTrue(result["output"].exists())
        self.assertTrue(result["qa_report"].exists())
        self.assertIn("heat_map_slide", result["patterns"])
        metadata = json.loads(result["manifest"].read_text(encoding="utf-8"))
        self.assertEqual(len(metadata["workflow_runs"]), 1)
        recorded = metadata["workflow_runs"][0]
        self.assertEqual(recorded["workflow"], "docx_to_pptx")
        self.assertEqual(recorded["palette"], "blackrock")
        self.assertIn("outputs\\model_brief_deck.pptx", recorded["outputs"])
        self.assertTrue((project_root / "treasury_decision" / "working" / "user_instructions.txt").exists())

    def test_unified_workflow_cli_reports_json_output(self):
        from docx import Document
        project_root = self.tmp / "projects"
        path = self.tmp / "brief.docx"
        doc = Document()
        doc.add_heading("Brief", 1)
        doc.add_paragraph("A clear recommendation.")
        doc.save(path)
        result = run([
            "artifact_workflow.py",
            "--source", str(path),
            "--project", "CLI Review",
            "--projects-root", str(project_root),
            "--json",
        ])
        self.assertEqual(result.returncode, 0, result.stderr)
        output = json.loads(result.stdout)
        self.assertEqual(output["run"]["workflow"], "docx_to_pptx")
        self.assertTrue(Path(output["output"]).exists())

    def test_pmse_workbook_pattern_saves_with_qa(self):
        from xlsx_patterns import pmse_comparison_workbook
        output = self.tmp / "pmse.xlsx"
        report = self.tmp / "pmse.qa.json"
        b = pmse_comparison_workbook(
            "PMSE Results",
            ["Bills", "Bonds"],
            ["Current", "Model B"],
            [[0.602, 0.141], [0.033, 0.004]],
        )
        b.save(output, final=True, report_path=report)
        self.assertTrue(output.exists())
        self.assertTrue(json.loads(report.read_text(encoding="utf-8"))["passed"])


class TestProjectWorkspace(unittest.TestCase):
    """Tests for grouped project input/output directory management."""

    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp())

    def tearDown(self):
        shutil.rmtree(self.tmp)

    def test_creates_expected_project_folders(self):
        from project_workspace import ProjectWorkspace
        workspace = ProjectWorkspace("Model Validation 2026", root=self.tmp)
        self.assertEqual(workspace.slug, "model_validation_2026")
        for folder in ["inputs", "assets", "working", "outputs", "qa"]:
            self.assertTrue((workspace.base_dir / folder).is_dir())
        metadata = json.loads((workspace.base_dir / "project.json").read_text(encoding="utf-8"))
        self.assertEqual(metadata["slug"], "model_validation_2026")
        self.assertEqual(metadata["workflow_runs"], [])

    def test_records_a_workflow_without_losing_folder_contract(self):
        from project_workspace import ProjectWorkspace
        workspace = ProjectWorkspace("Manifest Review", root=self.tmp)
        output = workspace.output_path("review.pptx")
        output.write_text("placeholder", encoding="utf-8")
        workspace.record_workflow(
            "docx_to_pptx",
            outputs=[output],
            palette="blackrock",
            patterns=["title_slide", "heat_map_slide"],
        )
        metadata = json.loads((workspace.base_dir / "project.json").read_text(encoding="utf-8"))
        self.assertEqual(metadata["folders"]["outputs"], "outputs")
        self.assertEqual(metadata["workflow_runs"][0]["patterns"], ["title_slide", "heat_map_slide"])

    def test_slide_plan_records_pattern_rationale_as_working_artifact(self):
        from decision_log import DecisionLog
        from project_workspace import ProjectWorkspace
        workspace = ProjectWorkspace("Agent Plan", root=self.tmp)
        log = DecisionLog.create(workspace.working_path("decision_log.md"))
        log.record_slide_plan([{
            "section": "Recommendation",
            "source": "Executive Summary",
            "assertion": "Approve Model B",
            "pattern": "recommendation_slide",
            "rationale": "A clear decision is required",
            "inputs": "Recommendation; supporting facts",
            "inference": "No",
        }])
        plan_path = workspace.working_path("slide_plan.md")
        self.assertTrue(plan_path.exists())
        plan = plan_path.read_text(encoding="utf-8")
        self.assertIn("`recommendation_slide`", plan)
        self.assertIn("A clear decision is required", plan)
        decision_log = workspace.working_path("decision_log.md").read_text(encoding="utf-8")
        self.assertIn("# Slide Plan", decision_log)


class TestArtifactUtilities(unittest.TestCase):
    """Tests for supported integration of the retained artifact utility bundle."""

    def setUp(self):
        self.tmp = Path(tempfile.mkdtemp())
        from project_workspace import ProjectWorkspace
        self.workspace = ProjectWorkspace("Advanced Utilities", root=self.tmp / "projects")

    def tearDown(self):
        shutil.rmtree(self.tmp)

    def test_catalog_exposes_every_retained_python_utility(self):
        from artifact_utilities import ArtifactUtilities, BUNDLE_ROOT
        tools = ArtifactUtilities(project=self.workspace)
        expected = len(list(BUNDLE_ROOT.rglob("*.py")))
        self.assertEqual(len(tools.capabilities), expected)
        self.assertIn("docx.a11y_audit", tools.capabilities)
        self.assertIn("slides.render_slides", tools.capabilities)
        self.assertIn("spreadsheets.spreadsheet_artifact_tool_starter", tools.capabilities)

    def test_optional_runtime_status_is_exposed(self):
        from artifact_utilities import ArtifactUtilities
        tools = ArtifactUtilities(project=self.workspace)
        records = {item["tool_id"]: item for item in tools.list()}
        self.assertIn("pdf2image", records["slides.render_slides"]["runtime_status"])
        self.assertIn("artifact_tool", records["spreadsheets.spreadsheet_artifact_tool_starter"]["runtime_status"])

    def test_exports_docx_table_to_project_working_folder(self):
        from artifact_utilities import ArtifactUtilities
        from docx import Document
        input_path = self.workspace.input_path("table.docx")
        doc = Document()
        table = doc.add_table(rows=2, cols=2)
        table.cell(0, 0).text = "Cohort"
        table.cell(0, 1).text = "PMSE"
        table.cell(1, 0).text = "Bills"
        table.cell(1, 1).text = "14.1%"
        doc.save(input_path)
        tools = ArtifactUtilities(project=self.workspace)
        output = tools.export_docx_table(input_path, "pmse_table.csv")
        self.assertTrue(output.exists())
        self.assertIn("Bills", output.read_text(encoding="utf-8"))

    def test_runs_docx_privacy_scrub_through_facade(self):
        from artifact_utilities import ArtifactUtilities
        from docx import Document
        input_path = self.workspace.input_path("private_note.docx")
        output_path = self.workspace.output_path("clean_note.docx")
        doc = Document()
        doc.core_properties.author = "Analyst Name"
        doc.add_paragraph("Shareable result.")
        doc.save(input_path)
        tools = ArtifactUtilities(project=self.workspace)
        tools.run("docx.privacy_scrub", [input_path.resolve(), "--out", output_path.resolve()])
        cleaned = Document(output_path)
        self.assertEqual(cleaned.core_properties.author or "", "")
        self.assertTrue((self.workspace.qa_dir / "docx_privacy_scrub.run.json").exists())

    def test_converts_xlsx_table_to_project_docx(self):
        from artifact_utilities import ArtifactUtilities
        from openpyxl import Workbook
        input_path = self.workspace.input_path("results.xlsx")
        workbook = Workbook()
        sheet = workbook.active
        sheet.title = "Results"
        sheet.append(["Cohort", "PMSE"])
        sheet.append(["Bills", 0.141])
        workbook.save(input_path)
        tools = ArtifactUtilities(project=self.workspace)
        output = tools.xlsx_table_to_docx(input_path, "results_table.docx", sheet="Results")
        self.assertTrue(output.exists())

    def test_creates_slide_montage_in_project_qa_folder(self):
        from PIL import Image
        from artifact_utilities import ArtifactUtilities
        slides_dir = self.workspace.working_path("slides")
        slides_dir.mkdir()
        Image.new("RGB", (320, 180), "#FFFFFF").save(slides_dir / "slide-1.png")
        Image.new("RGB", (320, 180), "#EEEEEE").save(slides_dir / "slide-2.png")
        tools = ArtifactUtilities(project=self.workspace)
        output = tools.create_montage(slides_dir)
        self.assertTrue(output.exists())

    def test_rendered_review_records_missing_optional_runtime(self):
        from unittest.mock import patch
        from artifact_utilities import ArtifactUtilities
        from create_pptx import PptxBuilder
        source = self.workspace.output_path("review.pptx")
        PptxBuilder().title_slide("Render Review").save(source, validate=False)
        tools = ArtifactUtilities(project=self.workspace)
        with patch("artifact_utilities.runtime_status", return_value={"soffice": False}):
            report_path = tools.review_rendered_pptx(source)
        report = json.loads(report_path.read_text(encoding="utf-8"))
        self.assertEqual(report["status"], "skipped")
        self.assertIn("soffice", report["missing_requirements"])


class TestComment(unittest.TestCase):
    """Tests for comment.py"""

    def setUp(self):
        self.tmp = tempfile.mkdtemp()
        from docx import Document
        path = Path(self.tmp) / "test.docx"
        doc = Document()
        doc.add_paragraph("Review this.")
        doc.save(path)
        self.out_dir = Path(self.tmp) / "unpacked"
        run(["unpack.py", str(path), str(self.out_dir)])

    def tearDown(self):
        shutil.rmtree(self.tmp)

    def test_add_comment(self):
        r = run(["comment.py", str(self.out_dir), "0", "Test comment"])
        self.assertEqual(r.returncode, 0, f"stdout: {r.stdout}")

    def test_comment_xml_created(self):
        run(["comment.py", str(self.out_dir), "0", "Hello"])
        self.assertTrue((self.out_dir / "word" / "comments.xml").exists())

    def test_comment_text_in_xml(self):
        run(["comment.py", str(self.out_dir), "0", "Review carefully"])
        content = (self.out_dir / "word" / "comments.xml").read_text(encoding="utf-8")
        self.assertIn("Review carefully", content)

    def test_custom_author(self):
        run(["comment.py", str(self.out_dir), "0", "Note", "--author", "Shikhar"])
        content = (self.out_dir / "word" / "comments.xml").read_text(encoding="utf-8")
        self.assertIn("Shikhar", content)


# ─── runner ───────────────────────────────────────────────────────────────────

if __name__ == "__main__":
    print("=" * 60)
    print("office_tools_mini test suite (Python-only)")
    print("=" * 60)

    missing = []
    for lib, pkg in [("docx", "python-docx"), ("pptx", "python-pptx"),
                     ("openpyxl", "openpyxl"), ("pandas", "pandas"),
                     ("PIL", "Pillow"), ("defusedxml", "defusedxml")]:
        try:
            __import__(lib)
        except ImportError:
            missing.append(pkg)

    if missing:
        print(f"\n⚠  Missing: {', '.join(missing)}")
        print("   Run: pip install -r requirements.txt\n")
    else:
        print("\n[OK] All dependencies found\n")

    unittest.main(verbosity=2)
