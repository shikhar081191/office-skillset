"""
Nineteen reusable PowerPoint slide patterns built on PptxBuilder.
Default palette: blackrock; supply any registered or external palette dictionary.

Each function appends one slide to a PptxBuilder and returns it.
Pass an existing builder to chain patterns into one deck, or let
the function create a new builder to get a standalone single-slide file.
"""

from pathlib import Path
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.oxml.ns import qn
from lxml import etree

from create_pptx import PALETTES, PptxBuilder, hex_to_rgb, resolve_palette

# ── default palette ───────────────────────────────────────────────────────────
BLACKROCK = resolve_palette(PALETTES["blackrock"])

RAG = {"green": "positive", "amber": "warning", "red": "negative"}


# ── shared helpers ─────────────────────────────────────────────────────────────

def _b(builder, palette):
    if builder is None:
        return PptxBuilder(palette=palette or BLACKROCK)
    if palette is not None:
        builder.palette = resolve_palette(palette)
    return builder


def _palette_builder(builder=None, palette=None):
    b = _b(builder, palette)
    return b, b.palette


def _mix(start, end, ratio):
    """Blend two palette colours and return a six-character RGB value."""
    ratio = max(0.0, min(float(ratio), 1.0))
    left = start.lstrip("#")
    right = end.lstrip("#")
    values = [
        round(int(left[i:i + 2], 16) * (1 - ratio) + int(right[i:i + 2], 16) * ratio)
        for i in (0, 2, 4)
    ]
    return "".join(f"{value:02X}" for value in values)


def _require_lengths(label, expected, actual):
    if expected != actual:
        raise ValueError(f"{label} length must be {expected}; received {actual}")


def _require_unit_interval(label, value):
    value = float(value)
    if not 0.0 <= value <= 1.0:
        raise ValueError(f"{label} must be between 0 and 1; received {value}")
    return value


def _format_value(value):
    value = float(value)
    if value.is_integer():
        return f"{value:,.0f}"
    return f"{value:,.2f}".rstrip("0").rstrip(".")


def _outlined_rectangle(slide, x, y, w, h, color, width=1.5, rounded=False):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE
        if rounded else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shape = slide.shapes.add_shape(shape_type, int(x), int(y), int(w), int(h))
    shape.fill.background()
    shape.line.color.rgb = hex_to_rgb(color)
    shape.line.width = Pt(width)
    return shape


def _header(slide, b, p, title, message=None):
    """Draw the standard slide header. Returns the y-inch offset where content should start."""
    b._fill_background(slide, p["background"])
    b._add_text_box(slide, title,
                    x=Inches(0.4), y=Inches(0.10), w=Inches(12.5), h=Inches(0.62),
                    font_size=22, bold=True, color=p["primary"])
    b._add_rectangle(slide,
                     x=Inches(0.4), y=Inches(0.75), w=Inches(12.5), h=Inches(0.04),
                     fill_color=p["secondary"])
    if message:
        b._add_text_box(slide, message,
                        x=Inches(0.4), y=Inches(0.82), w=Inches(12.5), h=Inches(0.34),
                        font_size=13, bold=True, color=p["primary"])
        return 1.20
    return 0.85


def _oval(slide, x, y, size, color):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL, int(x), int(y), int(size), int(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_to_rgb(color)
    shape.line.fill.background()
    return shape


def _set_cell_fill(cell, hex_color):
    """Set table cell background using direct XML (works across python-pptx versions)."""
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"):
        for elem in list(tcPr.findall(qn(tag))):
            tcPr.remove(elem)
    solid = etree.SubElement(tcPr, qn("a:solidFill"))
    clr = etree.SubElement(solid, qn("a:srgbClr"))
    clr.set("val", hex_color.lstrip("#").upper())


def _cell_text(cell, text, font_size=11, bold=False, color=None, align=PP_ALIGN.LEFT):
    tf = cell.text_frame
    tf.word_wrap = True
    para = tf.paragraphs[0]
    para.alignment = align
    for r in list(para._p.findall(qn("a:r"))):
        para._p.remove(r)
    run = para.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = hex_to_rgb(color or BLACKROCK["text_dark"])


def _looks_numeric(s):
    stripped = (
        s.strip()
        .lstrip("$€£+-")
        .rstrip("%x")
        .replace(",", "")
        .replace("bps", "")
    )
    try:
        float(stripped)
        return True
    except ValueError:
        return False


def _set_chart_series_colors(chart, colors):
    """Set fill color for each series (column/bar charts) via XML."""
    for i, series in enumerate(chart.series):
        if i >= len(colors):
            break
        ser_el = series._element
        spPr = ser_el.find(qn("c:spPr"))
        if spPr is None:
            spPr = etree.SubElement(ser_el, qn("c:spPr"))
        for tag in ("a:solidFill", "a:noFill", "a:gradFill", "a:pattFill"):
            for el in list(spPr.findall(qn(tag))):
                spPr.remove(el)
        sf = etree.SubElement(spPr, qn("a:solidFill"))
        clr = etree.SubElement(sf, qn("a:srgbClr"))
        clr.set("val", colors[i].lstrip("#").upper())


def _set_chart_line_colors(chart, colors):
    """Set stroke color for each series (line charts) via XML."""
    for i, series in enumerate(chart.series):
        if i >= len(colors):
            break
        ser_el = series._element
        spPr = ser_el.find(qn("c:spPr"))
        if spPr is None:
            spPr = etree.SubElement(ser_el, qn("c:spPr"))
        ln = spPr.find(qn("a:ln"))
        if ln is None:
            ln = etree.SubElement(spPr, qn("a:ln"))
        for tag in ("a:solidFill", "a:noFill"):
            for el in list(ln.findall(qn(tag))):
                ln.remove(el)
        sf = etree.SubElement(ln, qn("a:solidFill"))
        clr = etree.SubElement(sf, qn("a:srgbClr"))
        clr.set("val", colors[i].lstrip("#").upper())


def _lighten_chart_gridlines(chart, hex_color="EBEBEB"):
    """Set major gridlines to a very light colour (or remove if hex_color is None)."""
    val_ax = chart.value_axis._element
    mg = val_ax.find(qn("c:majorGridlines"))
    if mg is None:
        return
    spPr = mg.find(qn("c:spPr"))
    if spPr is None:
        spPr = etree.SubElement(mg, qn("c:spPr"))
    ln = spPr.find(qn("a:ln"))
    if ln is None:
        ln = etree.SubElement(spPr, qn("a:ln"))
    for child in list(ln):
        ln.remove(child)
    if hex_color:
        sf = etree.SubElement(ln, qn("a:solidFill"))
        clr = etree.SubElement(sf, qn("a:srgbClr"))
        clr.set("val", hex_color.lstrip("#").upper())
    else:
        ln.append(etree.Element(qn("a:noFill")))


# ── 1. Status ─────────────────────────────────────────────────────────────────

def status_slide(title, workstreams, key_message=None, builder=None, palette=None):
    """
    2×N grid of workstream status cards.

    workstreams: list of dicts — {name, status, rag: green|amber|red, owner}
    """
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    content_y = _header(slide, b, p, title, key_message)

    n_rows = (len(workstreams) + 1) // 2
    card_w_in = 5.9
    card_h_in = 1.7
    col_gap_in = 0.63
    row_gap_in = 0.15 if n_rows >= 4 else 0.22
    ox, oy = 0.4, content_y

    for i, ws in enumerate(workstreams):
        cx, cy = i % 2, i // 2
        x = Inches(ox + cx * (card_w_in + col_gap_in))
        y = Inches(oy + cy * (card_h_in + row_gap_in))
        cw, ch = Inches(card_w_in), Inches(card_h_in)

        b._add_rectangle(slide, x=x, y=y, w=cw, h=ch,
                         fill_color=p["surface"], line_color=p["border"])
        b._add_rectangle(slide, x=x, y=y, w=Inches(0.07), h=ch,
                         fill_color=p["secondary"])

        rag_color = p[RAG.get((ws.get("rag") or "green").lower(), RAG["green"])]
        dot_sz = Inches(0.28)
        dot_x = x + cw - dot_sz - Inches(0.18)
        dot_y = y + Inches(0.18)
        _oval(slide, dot_x, dot_y, dot_sz, rag_color)

        owner = str(ws.get("owner", ""))
        if owner:
            b._add_text_box(slide, owner,
                            x=dot_x - Inches(0.65), y=dot_y,
                            w=Inches(0.6), h=Inches(0.32),
                            font_size=11, bold=True, color=p["text_dark"],
                            align=PP_ALIGN.RIGHT)

        b._add_text_box(slide, ws.get("name", ""),
                        x=x + Inches(0.14), y=y + Inches(0.12),
                        w=Inches(card_w_in - 1.45), h=Inches(0.50),
                        font_size=15, bold=True, color=p["primary"])

        b._add_text_box(slide, ws.get("status", ""),
                        x=x + Inches(0.14), y=y + Inches(0.65),
                        w=Inches(card_w_in - 0.28), h=Inches(0.90),
                        font_size=11, color=p["muted_text"], wrap=True)

    return b


# ── 2. Process ────────────────────────────────────────────────────────────────

def process_slide(title, steps, highlight=None, footnote=None, key_message=None, builder=None, palette=None):
    """
    Horizontal strip of 3-5 boxes connected by arrows.

    steps:     list of dicts — {name, description}
    highlight: 0-based index of the step to colour in accent (optional)
    footnote:  string shown in small italic below the strip (optional)
    """
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    _header(slide, b, p, title, key_message)

    n = len(steps)
    arrow_w_in = 0.45
    usable_w_in = 12.5
    box_w_in = (usable_w_in - arrow_w_in * (n - 1)) / n
    box_h_in = 3.6
    box_y_in = 1.75
    start_x_in = 0.4

    for i, step in enumerate(steps):
        x_in = start_x_in + i * (box_w_in + arrow_w_in)
        x, y = Inches(x_in), Inches(box_y_in)
        bw, bh = Inches(box_w_in), Inches(box_h_in)
        is_hi = (i == highlight)

        box_fill = p["primary"] if is_hi else p["surface"]
        text_col = p["text_light"] if is_hi else p["text_dark"]
        num_col  = p["accent"]    if is_hi else p["primary"]
        desc_col = p["surface_alt"] if is_hi else p["muted_text"]

        b._add_rectangle(slide, x=x, y=y, w=bw, h=bh,
                         fill_color=box_fill, line_color=p["border"])

        b._add_text_box(slide, str(i + 1),
                        x=x, y=y + Inches(0.2), w=bw, h=Inches(0.8),
                        font_size=30, bold=True, color=num_col, align=PP_ALIGN.CENTER)

        b._add_text_box(slide, step.get("name", ""),
                        x=x + Inches(0.1), y=y + Inches(1.1),
                        w=Inches(box_w_in - 0.2), h=Inches(0.55),
                        font_size=11, bold=True, color=text_col, align=PP_ALIGN.CENTER)

        b._add_text_box(slide, step.get("description", ""),
                        x=x + Inches(0.1), y=y + Inches(1.75),
                        w=Inches(box_w_in - 0.2), h=Inches(1.7),
                        font_size=9, color=desc_col, align=PP_ALIGN.CENTER, wrap=True)

        if i < n - 1:
            ax = int(Inches(x_in + box_w_in))
            ay = int(Inches(box_y_in + box_h_in / 2 - 0.15))
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RIGHT_ARROW,
                ax, ay, int(Inches(arrow_w_in)), int(Inches(0.3))
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = hex_to_rgb(p["secondary"])
            arrow.line.fill.background()

    if footnote:
        b._add_text_box(slide, footnote,
                        x=Inches(0.4), y=Inches(6.85), w=Inches(12.5), h=Inches(0.45),
                        font_size=8, italic=True, color=p["muted_text"])

    return b


# ── 3. Timeline ───────────────────────────────────────────────────────────────

def timeline_slide(title, milestones, key_message=None, builder=None, palette=None):
    """
    Horizontal timeline with alternating above/below milestone labels.

    milestones: list of dicts — {date, label, state: past|current|future}
                Ordered earliest to latest.
    """
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    _header(slide, b, p, title, key_message)

    spine_x_in = 0.9
    spine_y_in = 3.85
    spine_w_in = 11.5
    spine_h_in = 0.06
    n = len(milestones)

    b._add_rectangle(slide,
                     x=Inches(spine_x_in), y=Inches(spine_y_in),
                     w=Inches(spine_w_in), h=Inches(spine_h_in),
                     fill_color=p["border"])

    spacing_in = spine_w_in / max(n - 1, 1) if n > 1 else 0.0

    for i, ms in enumerate(milestones):
        state = (ms.get("state") or "future").lower()
        tick_x_in = spine_x_in + i * spacing_in

        if state == "past":
            colour = p["muted_text"]
            dot_sz_in = 0.22
        elif state == "current":
            colour = p["secondary"]
            dot_sz_in = 0.38
        else:
            colour = p["accent"]
            dot_sz_in = 0.26

        tick_h_in = 0.42 if state == "current" else 0.28
        b._add_rectangle(slide,
                         x=Inches(tick_x_in - 0.03),
                         y=Inches(spine_y_in - tick_h_in / 2),
                         w=Inches(0.06), h=Inches(tick_h_in),
                         fill_color=colour)

        _oval(slide,
              Inches(tick_x_in - dot_sz_in / 2),
              Inches(spine_y_in + spine_h_in / 2 - dot_sz_in / 2),
              Inches(dot_sz_in), colour)

        if state == "current":
            ring_in = dot_sz_in + 0.18
            ring = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.OVAL,
                int(Inches(tick_x_in - ring_in / 2)),
                int(Inches(spine_y_in + spine_h_in / 2 - ring_in / 2)),
                int(Inches(ring_in)), int(Inches(ring_in))
            )
            ring.fill.background()
            ring.line.color.rgb = hex_to_rgb(p["secondary"])
            ring.line.width = Pt(2.0)

        above = (i % 2 == 0)
        lw_in = 2.2
        lx_in = min(max(tick_x_in - lw_in / 2, 0.4), 13.33 - 0.4 - lw_in)
        label_color = p["text_dark"]

        if above:
            label_y_in = spine_y_in - 1.95
            date_y_in  = spine_y_in - 1.18
            conn_y_in  = label_y_in + 0.62
            conn_h_in  = spine_y_in - conn_y_in
        else:
            label_y_in = spine_y_in + 0.60
            date_y_in  = spine_y_in + 1.30
            conn_y_in  = spine_y_in + spine_h_in
            conn_h_in  = label_y_in - conn_y_in

        b._add_text_box(slide, ms.get("label", ""),
                        x=Inches(lx_in), y=Inches(label_y_in),
                        w=Inches(lw_in), h=Inches(0.60),
                        font_size=14, bold=(state == "current"),
                        color=label_color, align=PP_ALIGN.CENTER)

        b._add_text_box(slide, ms.get("date", ""),
                        x=Inches(lx_in), y=Inches(date_y_in),
                        w=Inches(lw_in), h=Inches(0.44),
                        font_size=12, color=p["muted_text"], align=PP_ALIGN.CENTER)

        b._add_rectangle(slide,
                         x=Inches(tick_x_in - 0.01), y=Inches(conn_y_in),
                         w=Inches(0.02), h=Inches(conn_h_in),
                         fill_color=p["border"])

    return b


# ── 4. Results Comparison ─────────────────────────────────────────────────────

def results_slide(title, columns, rows, win_col=None, worst_col=None, key_message=None, builder=None, palette=None):
    """
    Structured comparison table with optional best/worst column highlighting.

    columns:   list of headers, e.g. ["Metric", "Model A", "Model B"]
    rows:      2-D list matching len(columns) per row
    win_col:   0-based column index to highlight green (optional)
    worst_col: 0-based column index to highlight amber (optional)
    """
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    b._fill_background(slide, p["background"])

    b._add_rectangle(slide, x=0, y=0, w=PptxBuilder.SLIDE_W, h=Inches(1.2),
                     fill_color=p["primary"])
    b._add_text_box(slide, title,
                    x=Inches(0.5), y=Inches(0.12), w=Inches(12.3), h=Inches(0.98),
                    font_size=28, bold=True, color=p["text_light"])

    tbl_y_in = 1.28
    if key_message:
        b._add_text_box(slide, key_message,
                        x=Inches(0.5), y=Inches(1.24), w=Inches(12.3), h=Inches(0.36),
                        font_size=13, bold=True, color=p["primary"])
        tbl_y_in = 1.66

    n_cols = len(columns)
    n_rows = len(rows) + 1
    tbl_x = Inches(0.5)
    tbl_y = Inches(tbl_y_in)
    tbl_w = Inches(12.3)
    tbl_h = Inches(7.5 - tbl_y_in - 0.18)

    gf = slide.shapes.add_table(n_rows, n_cols, tbl_x, tbl_y, tbl_w, tbl_h)
    tbl = gf.table

    label_w = int(int(tbl_w) * 0.20)
    data_w  = int((int(tbl_w) - label_w) / max(n_cols - 1, 1))
    tbl.columns[0].width = label_w
    for ci in range(1, n_cols):
        tbl.columns[ci].width = data_w

    for ci, hdr in enumerate(columns):
        cell = tbl.cell(0, ci)
        _set_cell_fill(cell, p["primary"])
        _cell_text(cell, hdr, font_size=14, bold=True, color=p["text_light"],
                   align=PP_ALIGN.CENTER)

    for ri, row in enumerate(rows):
        row_bg = p["background"] if ri % 2 == 0 else p["surface"]
        for ci, val in enumerate(row):
            cell = tbl.cell(ri + 1, ci)
            if win_col is not None and ci == win_col:
                _set_cell_fill(cell, _mix(p["background"], p["positive"], 0.18))
            elif worst_col is not None and ci == worst_col:
                _set_cell_fill(cell, p["highlight"])
            else:
                _set_cell_fill(cell, row_bg)

            is_number = ci > 0 and _looks_numeric(str(val))
            _cell_text(cell, val,
                       font_size=16 if is_number else 13,
                       bold=True,
                       color=p["text_dark"],
                       align=PP_ALIGN.CENTER)

    return b


# ── 5. Chart and Context ──────────────────────────────────────────────────────

def chart_context_slide(title, chart_path, headline, bullets, so_what, builder=None, palette=None):
    """
    Left 60% chart image, right 40% context panel.

    chart_path: file path to image (.png/.jpg). None or missing shows a placeholder.
    headline:   large bold summary in the right panel
    bullets:    list of 2-3 supporting points
    so_what:    single takeaway line shown in accent colour at the bottom
    """
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    _header(slide, b, p, title)

    chart_x = Inches(0.4)
    chart_y = Inches(0.82)
    chart_w = Inches(7.5)
    chart_h = Inches(6.3)

    cp = Path(chart_path) if chart_path else None
    if cp and cp.exists():
        slide.shapes.add_picture(str(cp), chart_x, chart_y, chart_w, chart_h)
    else:
        b._add_rectangle(slide, x=chart_x, y=chart_y, w=chart_w, h=chart_h,
                         fill_color=p["surface"], line_color=p["border"])
        b._add_text_box(slide, "[ chart ]",
                        x=chart_x, y=chart_y + Inches(2.7), w=chart_w, h=Inches(0.6),
                        font_size=16, color=p["muted_text"], align=PP_ALIGN.CENTER)

    div_x = chart_x + chart_w + Inches(0.18)
    b._add_rectangle(slide, x=div_x, y=chart_y, w=Inches(0.02), h=chart_h,
                     fill_color=p["border"])

    rx = div_x + Inches(0.22)
    rw = Inches(4.7)

    b._add_text_box(slide, headline,
                    x=rx, y=Inches(0.9), w=rw, h=Inches(1.3),
                    font_size=18, bold=True, color=p["primary"], wrap=True)

    bullets_text = "\n".join(f"•  {item}" for item in bullets)
    b._add_text_box(slide, bullets_text,
                    x=rx, y=Inches(2.35), w=rw, h=Inches(3.5),
                    font_size=11, color=p["text_dark"], wrap=True)

    b._add_rectangle(slide, x=rx, y=Inches(6.35), w=rw, h=Inches(0.6),
                     fill_color=p["accent"])
    b._add_text_box(slide, so_what,
                    x=rx + Inches(0.12), y=Inches(6.38),
                    w=rw - Inches(0.25), h=Inches(0.55),
                    font_size=10, bold=True, color=p["primary"], wrap=True)

    return b


# ── 6. Numbers ────────────────────────────────────────────────────────────────

def numbers_slide(title, stats, builder=None, palette=None):
    """
    3-4 large stat callouts on a dark background.

    stats: list of dicts — {number, label, context}
    """
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    b._fill_background(slide, p["primary"])

    b._add_text_box(slide, title,
                    x=Inches(0.8), y=Inches(0.3), w=Inches(11.7), h=Inches(0.7),
                    font_size=20, bold=True, color=p["text_light"])
    b._add_rectangle(slide,
                     x=Inches(0.8), y=Inches(1.05), w=Inches(11.7), h=Inches(0.05),
                     fill_color=p["secondary"])

    n = len(stats)
    box_w_in = 11.7 / n
    start_x_in = 0.8

    for i, stat in enumerate(stats):
        x_in = start_x_in + i * box_w_in

        if i < n - 1:
            b._add_rectangle(slide,
                             x=Inches(x_in + box_w_in - 0.01), y=Inches(1.4),
                             w=Inches(0.02), h=Inches(5.5),
                             fill_color=p["secondary"])

        b._add_text_box(slide, str(stat.get("number", "")),
                        x=Inches(x_in), y=Inches(2.0), w=Inches(box_w_in), h=Inches(2.2),
                        font_size=60, bold=True, color=p["accent"],
                        align=PP_ALIGN.CENTER)

        b._add_text_box(slide, str(stat.get("label", "")),
                        x=Inches(x_in), y=Inches(4.25), w=Inches(box_w_in), h=Inches(0.6),
                        font_size=15, color=p["text_light"], align=PP_ALIGN.CENTER)

        b._add_text_box(slide, str(stat.get("context", "")),
                        x=Inches(x_in + 0.1), y=Inches(4.95),
                        w=Inches(box_w_in - 0.2), h=Inches(0.9),
                        font_size=9, color=p["secondary"],
                        align=PP_ALIGN.CENTER, wrap=True)

    return b


# ── Demo ──────────────────────────────────────────────────────────────────────

def heat_map_slide(title, row_labels, col_labels, values, show_values=True,
                   key_message=None, builder=None, palette=None):
    """Matrix heat map coloured from the palette background to primary colour."""
    _require_lengths("values", len(row_labels), len(values))
    if not row_labels or not col_labels:
        raise ValueError("row_labels and col_labels must not be empty")
    for row in values:
        _require_lengths("heat-map row", len(col_labels), len(row))

    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    _header(slide, b, p, title, key_message)
    numeric = [[float(value) for value in row] for row in values]
    flattened = [value for row in numeric for value in row]
    minimum, maximum = min(flattened), max(flattened)
    spread = maximum - minimum
    x0, y0, grid_w, grid_h = 2.35, 1.4, 10.1, 5.25
    cell_w = grid_w / len(col_labels)
    cell_h = grid_h / len(row_labels)

    for ci, label in enumerate(col_labels):
        b._add_text_box(slide, str(label), x=Inches(x0 + ci * cell_w), y=Inches(0.92),
                        w=Inches(cell_w), h=Inches(0.38), font_size=9, bold=True,
                        color=p["primary"], align=PP_ALIGN.CENTER)
    for ri, label in enumerate(row_labels):
        y = y0 + ri * cell_h
        b._add_text_box(slide, str(label), x=Inches(0.42),
                        y=Inches(y + cell_h / 2 - 0.17), w=Inches(1.78), h=Inches(0.34),
                        font_size=9, bold=True, color=p["text_dark"])
        for ci, value in enumerate(numeric[ri]):
            x = x0 + ci * cell_w
            normalised = 0.5 if spread == 0 else (value - minimum) / spread
            fill = _mix(p["background"], p["primary"], 0.90 * normalised)
            b._add_rectangle(slide, x=Inches(x), y=Inches(y),
                             w=Inches(cell_w - 0.03), h=Inches(cell_h - 0.03),
                             fill_color=fill)
            if show_values:
                b._add_text_box(slide, _format_value(value), x=Inches(x),
                                y=Inches(y + cell_h / 2 - 0.17),
                                w=Inches(cell_w - 0.03), h=Inches(0.34),
                                font_size=10, bold=True,
                                color=p["text_light"] if normalised >= 0.58 else p["text_dark"],
                                align=PP_ALIGN.CENTER)
            if value == maximum:
                _outlined_rectangle(slide, Inches(x), Inches(y),
                                    Inches(cell_w - 0.03), Inches(cell_h - 0.03),
                                    p["accent"], width=2)
    return b


def waterfall_slide(title, items, total_label="Total", key_message=None, builder=None, palette=None):
    """Contribution bars against a zero line with a final calculated total."""
    if not items:
        raise ValueError("items must contain at least one contribution")
    parsed = [(str(item["name"]), float(item["value"])) for item in items]
    rows = parsed + [(str(total_label), sum(value for _, value in parsed))]
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    _header(slide, b, p, title, key_message)
    label_x, bar_w, spine_x = 0.5, 8.8, 7.05
    max_abs = max(abs(value) for _, value in rows) or 1.0
    unit_w = 3.7 / max_abs
    row_h = min(0.68, 5.6 / len(rows))
    y0 = 1.2
    b._add_rectangle(slide, x=Inches(spine_x), y=Inches(y0 - 0.08),
                     w=Inches(0.025), h=Inches(row_h * len(rows) + 0.15),
                     fill_color=p["border"])
    for index, (name, value) in enumerate(rows):
        y = y0 + index * row_h
        is_total = index == len(rows) - 1
        if is_total:
            b._add_rectangle(slide, x=Inches(label_x), y=Inches(y - 0.1),
                             w=Inches(bar_w + 2.2), h=Inches(0.02),
                             fill_color=p["border"])
        b._add_text_box(slide, name, x=Inches(label_x), y=Inches(y + 0.1),
                        w=Inches(2.35), h=Inches(0.3), font_size=10,
                        bold=is_total, color=p["text_dark"])
        length = abs(value) * unit_w
        x = spine_x if value >= 0 else spine_x - length
        fill = p["primary"] if is_total else (p["accent"] if value >= 0 else p["negative"])
        b._add_rectangle(slide, x=Inches(x), y=Inches(y + 0.08),
                         w=Inches(max(length, 0.03)), h=Inches(row_h - 0.22),
                         fill_color=fill)
        value_text = f"{'+' if value > 0 and not is_total else ''}{_format_value(value)}"
        tx = x + length + 0.12 if value >= 0 else x - 0.83
        b._add_text_box(slide, value_text, x=Inches(tx), y=Inches(y + 0.1),
                        w=Inches(0.75), h=Inches(0.3), font_size=9,
                        bold=is_total, color=p["text_dark"],
                        align=PP_ALIGN.LEFT if value >= 0 else PP_ALIGN.RIGHT)
    return b


def scorecard_slide(title, criteria, options, scores, weights=None,
                    key_message=None, builder=None, palette=None):
    """Model or vendor scorecard with numeric or RAG scoring and footer totals."""
    if not criteria or not options:
        raise ValueError("criteria and options must not be empty")
    _require_lengths("scores", len(criteria), len(scores))
    for row in scores:
        _require_lengths("scorecard row", len(options), len(row))
    if weights is not None:
        _require_lengths("weights", len(criteria), len(weights))
        if abs(sum(float(weight) for weight in weights) - 1.0) > 1e-6:
            raise ValueError("weights must sum to 1")
    is_numeric = all(isinstance(value, (int, float)) for row in scores for value in row)
    is_rag = all(str(value).lower() in RAG for row in scores for value in row)
    if not (is_numeric or is_rag):
        raise ValueError("scores must be either all numeric or all green/amber/red strings")
    if weights is not None and not is_numeric:
        raise ValueError("weights are supported only for numeric scorecards")

    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    _header(slide, b, p, title, key_message)
    rows, cols = len(criteria) + 2, len(options) + 1
    table = slide.shapes.add_table(rows, cols, Inches(0.5), Inches(1.05),
                                   Inches(12.25), Inches(5.95)).table
    table.columns[0].width = Inches(3.0)
    for ci in range(1, cols):
        table.columns[ci].width = int(Inches(9.25) / len(options))
    _set_cell_fill(table.cell(0, 0), p["primary"])
    _cell_text(table.cell(0, 0), "Criteria", bold=True, color=p["text_light"])
    for ci, option in enumerate(options, start=1):
        _set_cell_fill(table.cell(0, ci), p["primary"])
        _cell_text(table.cell(0, ci), option, bold=True, color=p["text_light"],
                   align=PP_ALIGN.CENTER)
    for ri, criterion in enumerate(criteria, start=1):
        _set_cell_fill(table.cell(ri, 0), p["surface"] if ri % 2 else p["background"])
        _cell_text(table.cell(ri, 0), criterion, bold=True, color=p["text_dark"])
        for ci, score in enumerate(scores[ri - 1], start=1):
            cell = table.cell(ri, ci)
            if is_numeric:
                value = float(score)
                fill = (_mix(p["background"], p["positive"], 0.22) if value >= 4
                        else _mix(p["background"], p["warning"], 0.25) if value <= 2
                        else p["surface"])
                _set_cell_fill(cell, fill)
                _cell_text(cell, _format_value(value), font_size=12, bold=True,
                           color=p["text_dark"], align=PP_ALIGN.CENTER)
            else:
                _set_cell_fill(cell, p["background"])
                _cell_text(cell, "●", font_size=18, bold=True,
                           color=p[RAG[str(score).lower()]], align=PP_ALIGN.CENTER)
    footer_row = len(criteria) + 1
    _set_cell_fill(table.cell(footer_row, 0), p["primary"])
    _cell_text(table.cell(footer_row, 0),
               "Weighted total" if weights is not None else ("Average" if is_numeric else "Overall"),
               bold=True, color=p["text_light"])
    for ci in range(len(options)):
        column = [scores[ri][ci] for ri in range(len(criteria))]
        if is_numeric:
            aggregate = (sum(float(value) * float(weights[ri])
                             for ri, value in enumerate(column))
                         if weights is not None
                         else sum(float(value) for value in column) / len(column))
            display = _format_value(aggregate)
            color = p["text_light"]
        else:
            overall = "red" if "red" in [str(v).lower() for v in column] else (
                "amber" if "amber" in [str(v).lower() for v in column] else "green"
            )
            display = "●"
            color = p[RAG[overall]]
        _set_cell_fill(table.cell(footer_row, ci + 1), p["primary"])
        _cell_text(table.cell(footer_row, ci + 1), display, font_size=12,
                   bold=True, color=color, align=PP_ALIGN.CENTER)
    return b


def annotation_chart_slide(title, chart_path, annotations, headline=None,
                           key_message=None, builder=None, palette=None):
    """Large chart image with labelled callouts located by relative coordinates."""
    chart = Path(chart_path)
    if not chart.exists():
        raise FileNotFoundError(f"Chart image not found: {chart_path}")
    for annotation in annotations:
        _require_unit_interval("annotation x", annotation["x"])
        _require_unit_interval("annotation y", annotation["y"])
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    _header(slide, b, p, title, key_message)
    chart_x, chart_y, chart_w, chart_h = 0.45, 0.9, 9.25, 6.1
    slide.shapes.add_picture(str(chart), Inches(chart_x), Inches(chart_y),
                             Inches(chart_w), Inches(chart_h))
    for index, annotation in enumerate(annotations):
        px = chart_x + float(annotation["x"]) * chart_w
        py = chart_y + float(annotation["y"]) * chart_h
        bx = min(max(px + 0.08, chart_x + 0.12), chart_x + chart_w - 2.35)
        by = max(chart_y + 0.1,
                 min(py - 0.55 - 0.14 * (index % 2), chart_y + chart_h - 0.55))
        connector = slide.shapes.add_connector(
            MSO_CONNECTOR.STRAIGHT, Inches(px), Inches(py), Inches(bx), Inches(by + 0.28)
        )
        connector.line.color.rgb = hex_to_rgb(p["accent"])
        connector.line.width = Pt(1.5)
        b._add_rectangle(slide, x=Inches(bx), y=Inches(by), w=Inches(2.25),
                         h=Inches(0.48), fill_color=p["background"],
                         line_color=p["accent"])
        b._add_text_box(slide, str(annotation["label"]), x=Inches(bx + 0.08),
                        y=Inches(by + 0.08), w=Inches(2.08), h=Inches(0.32),
                        font_size=8, bold=True, color=p["text_dark"])
    if headline:
        b._add_rectangle(slide, x=Inches(9.95), y=Inches(5.35),
                         w=Inches(2.85), h=Inches(1.45), fill_color=p["primary"])
        b._add_text_box(slide, str(headline), x=Inches(10.12), y=Inches(5.55),
                        w=Inches(2.5), h=Inches(1.05), font_size=12, bold=True,
                        color=p["text_light"], wrap=True)
    return b


def two_by_two_slide(title, x_label, y_label, x_low, x_high, y_low, y_high,
                     items, key_message=None, builder=None, palette=None):
    """Prioritisation matrix placing labelled items on normalised x/y axes."""
    for item in items:
        _require_unit_interval("item x", item["x"])
        _require_unit_interval("item y", item["y"])
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    content_y = _header(slide, b, p, title, key_message)
    left, top, width = 1.65, content_y + 0.06, 10.35
    height = 7.5 - top - 0.80
    half_w, half_h = width / 2, height / 2
    for x, y, fill in [
        (left, top, p["surface"]),
        (left + half_w, top, _mix(p["background"], p["accent"], 0.16)),
        (left, top + half_h, _mix(p["background"], p["primary"], 0.17)),
        (left + half_w, top + half_h, p["surface_alt"]),
    ]:
        b._add_rectangle(slide, x=Inches(x), y=Inches(y), w=Inches(half_w),
                         h=Inches(half_h), fill_color=fill)
    b._add_rectangle(slide, x=Inches(left + half_w), y=Inches(top),
                     w=Inches(0.025), h=Inches(height), fill_color=p["border"])
    b._add_rectangle(slide, x=Inches(left), y=Inches(top + half_h),
                     w=Inches(width), h=Inches(0.025), fill_color=p["border"])
    for label, x, y, w, align in [
        (x_low, left, top + height + 0.07, 2.0, PP_ALIGN.LEFT),
        (x_high, left + width - 2.0, top + height + 0.07, 2.0, PP_ALIGN.RIGHT),
        (y_high, 0.32, top - 0.02, 1.10, PP_ALIGN.RIGHT),
        (y_low, 0.32, top + height - 0.28, 1.10, PP_ALIGN.RIGHT),
    ]:
        b._add_text_box(slide, str(label), x=Inches(x), y=Inches(y),
                        w=Inches(w), h=Inches(0.28), font_size=10,
                        color=p["muted_text"], align=align)
    b._add_text_box(slide, x_label, x=Inches(left + width / 2 - 1.6),
                    y=Inches(top + height + 0.07), w=Inches(3.2), h=Inches(0.35),
                    font_size=14, bold=True, color=p["primary"], align=PP_ALIGN.CENTER)
    b._add_text_box(slide, y_label, x=Inches(0.28), y=Inches(top + height / 2 - 0.28),
                    w=Inches(1.2), h=Inches(0.55), font_size=14,
                    bold=True, color=p["primary"], align=PP_ALIGN.CENTER)
    for item in items:
        x = max(left + 0.06, min(left + float(item["x"]) * width - 0.9,
                                 left + width - 1.9))
        y = max(top + 0.06, min(top + (1 - float(item["y"])) * height - 0.38,
                                top + height - 0.85))
        descriptor = str(item.get("descriptor", ""))
        box_h = 0.72 if descriptor else 0.50
        box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                     Inches(x), Inches(y), Inches(1.80), Inches(box_h))
        box.fill.solid()
        box.fill.fore_color.rgb = hex_to_rgb(p["background"])
        box.line.color.rgb = hex_to_rgb(p["secondary"])
        b._add_text_box(slide, str(item["name"]), x=Inches(x + 0.07),
                        y=Inches(y + 0.06), w=Inches(1.66), h=Inches(0.28),
                        font_size=12, bold=True, color=p["text_dark"])
        if descriptor:
            b._add_text_box(slide, descriptor, x=Inches(x + 0.07),
                            y=Inches(y + 0.36), w=Inches(1.66), h=Inches(0.28),
                            font_size=11, color=p["muted_text"])
    return b


def assertion_evidence_slide(title, assertion, evidence, key_message=None, builder=None, palette=None):
    """Large assertion panel paired with concise supporting evidence."""
    if not evidence:
        raise ValueError("evidence must contain at least one item")
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    b._fill_background(slide, p["background"])
    b._add_rectangle(slide, x=0, y=0, w=Inches(6.0), h=PptxBuilder.SLIDE_H,
                     fill_color=p["primary"])
    b._add_text_box(slide, title.upper(), x=Inches(0.42), y=Inches(0.25),
                    w=Inches(4.9), h=Inches(0.32), font_size=10, bold=True,
                    color=p["surface_alt"])
    b._add_text_box(slide, assertion, x=Inches(0.55), y=Inches(1.8),
                    w=Inches(4.9), h=Inches(3.2), font_size=34, bold=True,
                    color=p["text_light"], wrap=True)
    right_y = 0.55
    if key_message:
        b._add_text_box(slide, key_message, x=Inches(6.45), y=Inches(0.28),
                        w=Inches(6.5), h=Inches(0.60), font_size=13, bold=True,
                        color=p["primary"], wrap=True)
        right_y = 1.02
    b._add_text_box(slide, "EVIDENCE", x=Inches(6.45), y=Inches(right_y),
                    w=Inches(6.5), h=Inches(0.32), font_size=10, bold=True,
                    color=p["muted_text"])
    y = right_y + 0.48
    row_h = min(1.5, (7.5 - y - 0.2) / max(len(evidence[:4]), 1))
    for item in evidence[:4]:
        stat = str(item.get("stat", ""))
        b._add_rectangle(slide, x=Inches(6.48), y=Inches(y + 0.08),
                         w=Inches(0.07), h=Inches(0.72), fill_color=p["accent"])
        text_x, text_w = 6.75, 6.0
        if stat:
            b._add_text_box(slide, stat, x=Inches(6.72), y=Inches(y),
                            w=Inches(1.2), h=Inches(0.44), font_size=20,
                            bold=True, color=p["primary"])
            text_x, text_w = 8.08, 4.68
        b._add_text_box(slide, str(item["text"]), x=Inches(text_x),
                        y=Inches(y + 0.04), w=Inches(text_w), h=Inches(0.78),
                        font_size=13, color=p["text_dark"], wrap=True)
        y += row_h
    return b


def diverging_bar_slide(title, row_labels, left_label, right_label,
                        left_values, right_values, key_message=None, builder=None, palette=None):
    """Two-sided comparison bars anchored to a shared centre spine."""
    _require_lengths("left_values", len(row_labels), len(left_values))
    _require_lengths("right_values", len(row_labels), len(right_values))
    if not row_labels:
        raise ValueError("row_labels must not be empty")
    left = [float(value) for value in left_values]
    right = [float(value) for value in right_values]
    if any(value < 0 for value in left + right):
        raise ValueError("diverging bar values must be positive")
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    _header(slide, b, p, title, key_message)
    spine_x, max_width = 6.66, 4.5
    max_value = max(left + right) or 1.0
    row_h = min(0.72, 5.2 / len(row_labels))
    y0 = 1.48
    gaps = [abs(lval - rval) for lval, rval in zip(left, right)]
    max_gap = max(gaps)
    b._add_text_box(slide, left_label, x=Inches(1.25), y=Inches(0.95),
                    w=Inches(4.7), h=Inches(0.3), font_size=11, bold=True,
                    color=p["primary"], align=PP_ALIGN.RIGHT)
    b._add_text_box(slide, right_label, x=Inches(7.32), y=Inches(0.95),
                    w=Inches(4.7), h=Inches(0.3), font_size=11, bold=True,
                    color=p["primary"])
    for index, gap in enumerate(gaps):
        if gap == max_gap:
            y = y0 + index * row_h
            b._add_rectangle(slide, x=Inches(0.7), y=Inches(y - 0.04),
                             w=Inches(12.0), h=Inches(row_h - 0.05),
                             fill_color=p["highlight"])
    b._add_rectangle(slide, x=Inches(spine_x), y=Inches(1.28),
                     w=Inches(0.025), h=Inches(row_h * len(row_labels) + 0.18),
                     fill_color=p["primary"])
    for index, label in enumerate(row_labels):
        y = y0 + index * row_h
        lw, rw = left[index] / max_value * max_width, right[index] / max_value * max_width
        b._add_rectangle(slide, x=Inches(spine_x - lw), y=Inches(y + 0.12),
                         w=Inches(lw), h=Inches(row_h - 0.26),
                         fill_color=p["secondary"])
        b._add_rectangle(slide, x=Inches(spine_x + 0.025), y=Inches(y + 0.12),
                         w=Inches(rw), h=Inches(row_h - 0.26), fill_color=p["accent"])
        b._add_text_box(slide, _format_value(left[index]),
                        x=Inches(spine_x - lw - 0.72), y=Inches(y + 0.15),
                        w=Inches(0.65), h=Inches(0.24), font_size=8, bold=True,
                        color=p["text_dark"], align=PP_ALIGN.RIGHT)
        b._add_text_box(slide, label, x=Inches(spine_x - 0.7),
                        y=Inches(y + 0.38), w=Inches(1.42), h=Inches(0.22),
                        font_size=8, color=p["text_dark"], align=PP_ALIGN.CENTER)
        b._add_text_box(slide, _format_value(right[index]),
                        x=Inches(spine_x + rw + 0.12), y=Inches(y + 0.15),
                        w=Inches(0.65), h=Inches(0.24), font_size=8, bold=True,
                        color=p["text_dark"])
    return b


# ── 14. Recommendation ───────────────────────────────────────────────────────

def recommendation_slide(title, recommendation, rationale, caveats=None,
                         key_message=None, builder=None, palette=None):
    """Dark-panel decision slide.

    recommendation: concise decision statement (~15 words) shown large on the left
    rationale:      list of supporting reasons shown on the right
    caveats:        optional list of risk / caveat bullets below the rationale
    key_message:    one-sentence summary shown bold at the top of the right panel
    """
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    b._fill_background(slide, p["background"])
    b._add_rectangle(slide, x=0, y=0, w=Inches(5.9), h=PptxBuilder.SLIDE_H,
                     fill_color=p["primary"])
    b._add_text_box(slide, "RECOMMENDATION", x=Inches(0.42), y=Inches(0.28),
                    w=Inches(5.0), h=Inches(0.32), font_size=10, bold=True,
                    color=p["surface_alt"])
    b._add_text_box(slide, title, x=Inches(0.42), y=Inches(0.65),
                    w=Inches(5.0), h=Inches(0.32), font_size=12,
                    color=p["surface_alt"])
    b._add_text_box(slide, recommendation, x=Inches(0.48), y=Inches(1.45),
                    w=Inches(4.9), h=Inches(4.1), font_size=30, bold=True,
                    color=p["text_light"], wrap=True, auto_fit=True)

    right_y = 0.28
    if key_message:
        b._add_text_box(slide, key_message, x=Inches(6.25), y=Inches(right_y),
                        w=Inches(6.8), h=Inches(0.68), font_size=14, bold=True,
                        color=p["primary"], wrap=True)
        right_y = 1.05
    b._add_text_box(slide, "RATIONALE", x=Inches(6.25), y=Inches(right_y),
                    w=Inches(6.8), h=Inches(0.32), font_size=10, bold=True,
                    color=p["muted_text"])
    rationale_top = right_y + 0.38
    rationale_h_in = (7.5 - rationale_top - 0.2) if not caveats else 3.85
    body = "\n".join(f"•  {r}" for r in rationale)
    b._add_text_box(slide, body, x=Inches(6.25), y=Inches(rationale_top),
                    w=Inches(6.8), h=Inches(rationale_h_in), font_size=16, bold=True,
                    color=p["text_dark"], wrap=True, auto_fit=True)
    if caveats:
        caveat_y = rationale_top + rationale_h_in + 0.12
        b._add_rectangle(slide, x=Inches(6.25), y=Inches(caveat_y),
                         w=Inches(6.8), h=Inches(0.03), fill_color=p["border"])
        b._add_text_box(slide, "CAVEATS", x=Inches(6.25), y=Inches(caveat_y + 0.08),
                        w=Inches(6.8), h=Inches(0.32), font_size=10, bold=True,
                        color=p["negative"])
        caveat_text = "\n".join(f"•  {c}" for c in caveats)
        b._add_text_box(slide, caveat_text, x=Inches(6.25), y=Inches(caveat_y + 0.46),
                        w=Inches(6.8), h=Inches(1.5), font_size=13,
                        color=p["muted_text"], wrap=True, auto_fit=True)
    return b


# ── 15. KPI Dashboard ────────────────────────────────────────────────────────

def kpi_dashboard_slide(title, kpis, key_message=None, builder=None, palette=None):
    """Grid of up to 6 KPI metric cards.

    kpis: list of dicts — name (str), value (str), delta (str, optional),
          trend ('up' | 'down' | 'flat', optional)
    """
    import math as _math
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    content_y = _header(slide, b, p, title, key_message)

    n = min(len(kpis), 6)
    cols = 3 if n > 3 else max(n, 1)
    rows = _math.ceil(n / cols)
    card_w, card_h, gap = 3.9, 2.5, 0.25
    total_w = cols * card_w + (cols - 1) * gap
    x0 = (13.33 - total_w) / 2
    y0 = content_y + (7.5 - content_y - rows * card_h - (rows - 1) * gap) / 2

    for i, kpi in enumerate(kpis[:n]):
        col, row = i % cols, i // cols
        cx = x0 + col * (card_w + gap)
        cy = y0 + row * (card_h + gap)
        b._add_rectangle(slide, x=Inches(cx), y=Inches(cy),
                         w=Inches(card_w), h=Inches(card_h),
                         fill_color=p["surface"], line_color=p["border"])
        b._add_text_box(slide, kpi["name"].upper(),
                        x=Inches(cx + 0.2), y=Inches(cy + 0.18),
                        w=Inches(card_w - 0.4), h=Inches(0.36),
                        font_size=11, bold=True, color=p["muted_text"])
        b._add_text_box(slide, str(kpi["value"]),
                        x=Inches(cx + 0.2), y=Inches(cy + 0.58),
                        w=Inches(card_w - 0.4), h=Inches(1.1),
                        font_size=44, bold=True, color=p["primary"])
        delta = str(kpi.get("delta", ""))
        trend = str(kpi.get("trend", ""))
        if delta or trend:
            trend_color = (
                p["positive"] if trend == "up" else
                p["negative"] if trend == "down" else
                p["muted_text"]
            )
            arrow = "^ " if trend == "up" else "v " if trend == "down" else "- "
            b._add_text_box(slide, f"{arrow}{delta}",
                            x=Inches(cx + 0.2), y=Inches(cy + 1.82),
                            w=Inches(card_w - 0.4), h=Inches(0.42),
                            font_size=13, bold=True, color=trend_color)
    return b


# ── 16. Agenda ───────────────────────────────────────────────────────────────

def agenda_slide(title, items, key_message=None, builder=None, palette=None):
    """Numbered agenda list with optional current-item highlight.

    items: list of dicts — label (str), description (str, optional),
           active (bool, optional — marks the current agenda item)
    """
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    content_y = _header(slide, b, p, title, key_message)

    n = min(len(items), 8)
    row_h = min(0.95, (7.5 - content_y - 0.3) / max(n, 1))
    y0 = content_y + 0.08

    for i, item in enumerate(items[:n]):
        y = y0 + i * row_h
        is_active = bool(item.get("active", False))
        if is_active:
            b._add_rectangle(slide, x=Inches(0.5), y=Inches(y - 0.06),
                             w=Inches(12.3), h=Inches(row_h - 0.04),
                             fill_color=p["highlight"])
        b._add_rectangle(slide, x=Inches(0.55), y=Inches(y + 0.05),
                         w=Inches(0.52), h=Inches(0.52),
                         fill_color=p["primary"] if is_active else p["surface"])
        b._add_text_box(slide, str(i + 1),
                        x=Inches(0.55), y=Inches(y + 0.07),
                        w=Inches(0.52), h=Inches(0.40), font_size=13, bold=True,
                        color=p["text_light"] if is_active else p["muted_text"],
                        align=PP_ALIGN.CENTER)
        label_color = p["text_dark"] if is_active else p["muted_text"]
        b._add_text_box(slide, item["label"],
                        x=Inches(1.24), y=Inches(y + 0.04),
                        w=Inches(11.5), h=Inches(0.46),
                        font_size=18, bold=is_active, color=label_color)
        if item.get("description"):
            b._add_text_box(slide, item["description"],
                            x=Inches(1.24), y=Inches(y + 0.52),
                            w=Inches(11.5), h=Inches(0.36),
                            font_size=12, color=p["muted_text"])
    return b


# ── 17. Assumption Table ─────────────────────────────────────────────────────

def assumption_table_slide(title, assumptions, footnote=None,
                           key_message=None, builder=None, palette=None):
    """Structured model assumptions table.

    assumptions: list of dicts — assumption (str), value (str),
                 source (str, optional), sensitivity ('high'|'medium'|'low', optional)
    """
    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    content_y = _header(slide, b, p, title, key_message)

    has_source = any(a.get("source") for a in assumptions)
    has_sens = any(a.get("sensitivity") for a in assumptions)
    col_names = ["Assumption", "Value"]
    col_w_raw = [5.0, 2.8]
    if has_source:
        col_names.append("Source")
        col_w_raw.append(2.5)
    if has_sens:
        col_names.append("Sensitivity")
        col_w_raw.append(1.6)
    scale = 11.8 / sum(col_w_raw)
    col_widths = [w * scale for w in col_w_raw]

    n_data = min(len(assumptions), 12)
    n_rows = n_data + 1
    tbl_y_in = content_y + 0.08
    avail_h = 7.5 - tbl_y_in - 0.35
    row_h = min(Inches(0.62), Inches(avail_h / max(n_rows, 1)))
    tbl = slide.shapes.add_table(
        n_rows, len(col_names),
        Inches(0.75), Inches(tbl_y_in),
        Inches(11.8), row_h * n_rows,
    )
    table = tbl.table
    for ci, cw in enumerate(col_widths):
        table.columns[ci].width = Inches(cw)

    for ci, name in enumerate(col_names):
        cell = table.cell(0, ci)
        _set_cell_fill(cell, p["primary"])
        _cell_text(cell, name, font_size=12, bold=True, color=p["text_light"])

    sens_colors = {"high": p["negative"], "medium": p["warning"], "low": p["positive"]}
    for ri, a in enumerate(assumptions[:n_data]):
        row_bg = p["background"] if ri % 2 == 0 else p["surface"]
        values = [a.get("assumption", ""), a.get("value", "")]
        if has_source:
            values.append(a.get("source", ""))
        if has_sens:
            values.append(a.get("sensitivity", "").capitalize())
        for ci, val in enumerate(values):
            cell = table.cell(ri + 1, ci)
            _set_cell_fill(cell, row_bg)
            text_color = p["text_dark"]
            is_bold = ci == 0
            if has_sens and ci == len(col_names) - 1:
                sens_key = a.get("sensitivity", "").lower()
                if sens_key in sens_colors:
                    text_color = sens_colors[sens_key]
                    is_bold = True
            _cell_text(cell, val, font_size=13, bold=is_bold, color=text_color)

    if footnote:
        b._add_text_box(slide, footnote,
                        x=Inches(0.75), y=Inches(6.97),
                        w=Inches(11.8), h=Inches(0.28),
                        font_size=7, color=p["muted_text"], italic=True)
    return b


# ── 18. Bar Chart ────────────────────────────────────────────────────────────

def bar_chart_slide(title, categories, series, footnote=None,
                    key_message=None, builder=None, palette=None):
    """Native clustered column (bar) chart.

    categories: list of category label strings
    series:     list of dicts — name (str), values (list of float)
    """
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    content_y = _header(slide, b, p, title, key_message)

    cd = ChartData()
    cd.categories = [str(c) for c in categories]
    for s in series:
        cd.add_series(str(s["name"]), tuple(float(v) for v in s["values"]))

    chart_top = content_y + 0.10
    chart_h = 7.5 - chart_top - (0.55 if footnote else 0.18)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(0.75), Inches(chart_top), Inches(11.8), Inches(chart_h),
        cd,
    )
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.has_major_gridlines = False

    _set_chart_series_colors(chart, [p["secondary"], p["primary"], p["accent"],
                                     p["muted_text"], p["positive"], p["negative"]])
    _lighten_chart_gridlines(chart)

    if footnote:
        b._add_text_box(slide, footnote,
                        x=Inches(0.75), y=Inches(6.97),
                        w=Inches(11.8), h=Inches(0.28),
                        font_size=7, color=p["muted_text"], italic=True)
    return b


# ── 19. Line Chart ───────────────────────────────────────────────────────────

def line_chart_slide(title, categories, series, footnote=None,
                     key_message=None, builder=None, palette=None):
    """Native line chart with one or more series.

    categories: list of category label strings
    series:     list of dicts — name (str), values (list of float)
    """
    from pptx.chart.data import ChartData
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

    b, p = _palette_builder(builder, palette)
    slide = b.prs.slides.add_slide(b._blank_layout)
    content_y = _header(slide, b, p, title, key_message)

    cd = ChartData()
    cd.categories = [str(c) for c in categories]
    for s in series:
        cd.add_series(str(s["name"]), tuple(float(v) for v in s["values"]))

    chart_top = content_y + 0.10
    chart_h = 7.5 - chart_top - (0.55 if footnote else 0.18)
    chart_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(0.75), Inches(chart_top), Inches(11.8), Inches(chart_h),
        cd,
    )
    chart = chart_frame.chart
    chart.has_title = False
    chart.has_legend = len(series) > 1
    if chart.has_legend:
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False
    chart.value_axis.has_major_gridlines = True
    chart.category_axis.has_major_gridlines = False

    _set_chart_line_colors(chart, [p["muted_text"], p["secondary"], p["accent"],
                                   p["primary"], p["positive"], p["negative"]])
    _lighten_chart_gridlines(chart)

    if footnote:
        b._add_text_box(slide, footnote,
                        x=Inches(0.75), y=Inches(6.97),
                        w=Inches(11.8), h=Inches(0.28),
                        font_size=7, color=p["muted_text"], italic=True)
    return b


def _demo_chart_image(output_path):
    """Create a lightweight chart image for the annotation pattern demo."""
    from PIL import Image, ImageDraw

    p = BLACKROCK
    image = Image.new("RGB", (1100, 650), "#" + p["background"])
    draw = ImageDraw.Draw(image)
    draw.line((90, 70, 90, 560), fill="#" + p["border"], width=3)
    draw.line((90, 560, 1030, 560), fill="#" + p["border"], width=3)
    points = [(110, 430), (245, 390), (360, 405), (480, 260),
              (600, 225), (710, 325), (820, 170), (990, 205)]
    draw.line(points, fill="#" + p["secondary"], width=7, joint="curve")
    for point in points:
        draw.ellipse((point[0] - 8, point[1] - 8, point[0] + 8, point[1] + 8),
                     fill="#" + p["accent"], outline="#" + p["primary"], width=2)
    image.save(output_path)


def demo_all():
    Path("output").mkdir(exist_ok=True)

    status_slide(
        title="Programme Status — Q2 2025",
        workstreams=[
            {"name": "Data Infrastructure",  "status": "Pipeline migration complete; UAT in progress",     "rag": "green", "owner": "AK"},
            {"name": "Risk Models",           "status": "Factor model 2w behind — vendor data late",        "rag": "amber", "owner": "SP"},
            {"name": "Client Reporting",      "status": "MVP live; stakeholder feedback incorporated",      "rag": "green", "owner": "LM"},
            {"name": "Regulatory Compliance", "status": "MiFID II gap analysis blocked on legal sign-off",  "rag": "red",   "owner": "TJ"},
            {"name": "Technology Platform",   "status": "API refresh on track; performance tests passed",   "rag": "green", "owner": "RM"},
            {"name": "Talent & Resourcing",   "status": "Two open roles filled; onboarding complete",       "rag": "green", "owner": "CB"},
        ]
    ).save("output/demo_status.pptx")
    print("output/demo_status.pptx")

    process_slide(
        title="Trade Lifecycle — Front to Back",
        steps=[
            {"name": "Order Capture",  "description": "OMS receives and validates order from PM"},
            {"name": "Pre-Trade Risk", "description": "Real-time limit checks against portfolio constraints"},
            {"name": "Execution",      "description": "Smart-order routing to lit and dark venues"},
            {"name": "Confirmation",   "description": "Matched via SWIFT or FIX confirmation messages"},
            {"name": "Settlement",     "description": "T+2 DVP settlement via CSD or prime broker"},
        ],
        highlight=2,
        footnote="* All times are T+0 unless subject to a market-specific settlement convention.",
    ).save("output/demo_process.pptx")
    print("output/demo_process.pptx")

    timeline_slide(
        title="Platform Delivery Roadmap — 2025",
        milestones=[
            {"date": "Jan 25", "label": "Discovery",             "state": "past"},
            {"date": "Mar 25", "label": "Architecture Sign-off", "state": "past"},
            {"date": "May 25", "label": "MVP Build",             "state": "current"},
            {"date": "Jul 25", "label": "UAT",                   "state": "future"},
            {"date": "Sep 25", "label": "Go-Live",               "state": "future"},
            {"date": "Dec 25", "label": "Optimise & Embed",      "state": "future"},
        ]
    ).save("output/demo_timeline.pptx")
    print("output/demo_timeline.pptx")

    results_slide(
        title="Model Backtesting — Return & Risk Comparison",
        columns=["Metric", "Baseline", "Model A", "Model B", "Delta vs Base"],
        rows=[
            ["Annualised Return", "8.2%",  "9.4%",  "11.1%", "+2.9%"],
            ["Volatility",        "14.1%", "13.8%", "12.6%", "-1.5%"],
            ["Sharpe Ratio",      "0.58",  "0.68",  "0.88",  "+0.30"],
            ["Max Drawdown",      "-18%",  "-15%",  "-11%",  "+7%"],
            ["Tracking Error",    "n/a",   "2.1%",  "2.4%",  "—"],
        ],
        win_col=3,
        worst_col=1,
    ).save("output/demo_results.pptx")
    print("output/demo_results.pptx")

    chart_context_slide(
        title="Equity Vol Regime — Q2 2025",
        chart_path=None,
        headline="Volatility has normalised following the February spike",
        bullets=[
            "VIX mean-reverted to 14.2 — 3yr average is 16.1",
            "Realised vol (30d) tracking implied at 0.97 ratio",
            "Tail risk premium compressed 40bps quarter-on-quarter",
        ],
        so_what="Position for continued compression: reduce vol overlay by 20%",
    ).save("output/demo_chart_context.pptx")
    print("output/demo_chart_context.pptx")

    numbers_slide(
        title="Portfolio at a Glance — 31 May 2025",
        stats=[
            {"number": "$2.4B", "label": "AUM",           "context": "+12% vs prior year"},
            {"number": "0.91",  "label": "Sharpe Ratio",  "context": "Above 0.80 target"},
            {"number": "-11%",  "label": "Max Drawdown",  "context": "Within -15% limit"},
            {"number": "94bps", "label": "Active Return", "context": "vs benchmark Q2 YTD"},
        ]
    ).save("output/demo_numbers.pptx")
    print("output/demo_numbers.pptx")

    heat_map_slide(
        title="PMSE Error Surface by Model and Asset Cohort",
        row_labels=["Rates", "Credit", "Equities", "FX"],
        col_labels=["Baseline", "Model A", "Model B", "Ensemble"],
        values=[[0.62, 0.55, 0.48, 0.44],
                [0.71, 0.61, 0.52, 0.46],
                [0.81, 0.69, 0.58, 0.50],
                [0.49, 0.43, 0.39, 0.35]],
    ).save("output/demo_heat_map_slide.pptx")
    print("output/demo_heat_map_slide.pptx")

    waterfall_slide(
        title="Contribution to PMSE Improvement vs Baseline",
        items=[
            {"name": "Rates factors", "value": 12.5},
            {"name": "Credit spreads", "value": 8.2},
            {"name": "Volatility regime", "value": 5.1},
            {"name": "FX carry noise", "value": -2.8},
            {"name": "Residual effects", "value": -1.4},
        ],
        total_label="Net improvement (bps)",
    ).save("output/demo_waterfall_slide.pptx")
    print("output/demo_waterfall_slide.pptx")

    scorecard_slide(
        title="Candidate Model Selection Scorecard",
        criteria=["PMSE accuracy", "Stability", "Explainability", "Runtime", "Governance"],
        options=["Baseline", "Model B", "Ensemble"],
        scores=[[2, 4, 5], [4, 4, 4], [5, 4, 3], [5, 4, 2], [5, 4, 3]],
        weights=[0.35, 0.20, 0.15, 0.10, 0.20],
    ).save("output/demo_scorecard_slide.pptx")
    print("output/demo_scorecard_slide.pptx")

    chart_path = Path("output/demo_annotation_chart_source.png")
    _demo_chart_image(chart_path)
    annotation_chart_slide(
        title="Rolling Equity-Credit Correlation: Regime Signals",
        chart_path=chart_path,
        annotations=[
            {"label": "Liquidity shock", "x": 0.42, "y": 0.36},
            {"label": "Policy pivot", "x": 0.78, "y": 0.20},
        ],
        headline="Correlation risk re-entered the stress band",
    ).save("output/demo_annotation_chart_slide.pptx")
    print("output/demo_annotation_chart_slide.pptx")

    two_by_two_slide(
        title="Validation Priorities: Effort vs Risk Reduction",
        x_label="Implementation effort", y_label="Risk reduction",
        x_low="Low effort", x_high="High effort",
        y_low="Low impact", y_high="High impact",
        items=[
            {"name": "Drift alerts", "x": 0.22, "y": 0.78, "descriptor": "Quick win"},
            {"name": "DCC rebuild", "x": 0.81, "y": 0.86, "descriptor": "Strategic"},
            {"name": "Data lineage", "x": 0.48, "y": 0.56},
            {"name": "UI polish", "x": 0.30, "y": 0.20},
        ],
    ).save("output/demo_two_by_two_slide.pptx")
    print("output/demo_two_by_two_slide.pptx")

    assertion_evidence_slide(
        title="Core Finding",
        assertion="Bills drive almost all of the PMSE improvement",
        evidence=[
            {"stat": "46pp", "text": "lower PMSE in short-duration sovereign cohorts"},
            {"stat": "p<0.01", "text": "improvement remains significant after regime controls"},
            {"stat": "82%", "text": "of aggregate gain is concentrated in bill exposures"},
        ],
    ).save("output/demo_assertion_evidence_slide.pptx")
    print("output/demo_assertion_evidence_slide.pptx")

    diverging_bar_slide(
        title="PMSE Comparison by Asset Cohort",
        row_labels=["Bills", "UST 2Y", "UST 10Y", "IG Credit", "Equity Beta", "FX Carry"],
        left_label="Current Model", right_label="Model B",
        left_values=[0.74, 0.65, 0.59, 0.68, 0.77, 0.49],
        right_values=[0.36, 0.44, 0.48, 0.52, 0.61, 0.40],
    ).save("output/demo_diverging_bar_slide.pptx")
    print("output/demo_diverging_bar_slide.pptx")


if __name__ == "__main__":
    demo_all()
