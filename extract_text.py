"""Extract text content from Office files (DOCX, PPTX, XLSX).

Outputs text under section headers for easy reading and analysis.

Usage:
    python extract_text.py document.docx
    python extract_text.py presentation.pptx
    python extract_text.py spreadsheet.xlsx
    python extract_text.py file.xlsx --head 50
"""

import argparse
import sys
from pathlib import Path


def extract_text(input_file: str, head: int = None) -> str:
    path = Path(input_file)
    if not path.exists():
        return f"Error: {input_file} does not exist"

    suffix = path.suffix.lower()

    try:
        if suffix == ".docx":
            text = _extract_docx(path)
        elif suffix == ".pptx":
            text = _extract_pptx(path)
        elif suffix in {".xlsx", ".xlsm", ".csv", ".tsv"}:
            text = _extract_xlsx(path)
        else:
            return f"Error: Unsupported file type: {suffix}"

        if head:
            lines = text.splitlines()
            text = "\n".join(lines[:head])
            if len(lines) > head:
                text += f"\n... ({len(lines) - head} more lines)"

        return text

    except ImportError as e:
        missing = str(e).split("'")[1] if "'" in str(e) else str(e)
        return f"Error: Missing library. Run: pip install {missing}"
    except Exception as e:
        return f"Error extracting text: {e}"


def _extract_docx(path: Path) -> str:
    from docx import Document
    doc = Document(path)
    lines = []

    for para in doc.paragraphs:
        text = para.text.strip()
        if text:
            style = para.style.name if para.style else ""
            if "Heading" in style:
                level = style.replace("Heading ", "").strip()
                prefix = "#" * int(level) if level.isdigit() else "#"
                lines.append(f"{prefix} {text}")
            else:
                lines.append(text)

    # Extract tables
    for i, table in enumerate(doc.tables):
        lines.append(f"\n## Table {i + 1}")
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            lines.append(" | ".join(cells))

    return "\n".join(lines)


def _extract_pptx(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(path)
    lines = []

    for i, slide in enumerate(prs.slides):
        lines.append(f"\n## Slide {i + 1}")

        # Collect shapes sorted by vertical position
        shapes = sorted(
            [s for s in slide.shapes if s.has_text_frame],
            key=lambda s: s.top if s.top is not None else 0,
        )

        for shape in shapes:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    lines.append(text)

    return "\n".join(lines)


def _extract_xlsx(path: Path) -> str:
    try:
        import pandas as pd
        all_sheets = pd.read_excel(path, sheet_name=None, dtype=str)
        lines = []
        for sheet_name, df in all_sheets.items():
            lines.append(f"\n## Sheet: {sheet_name}")
            lines.append(df.fillna("").to_string(index=False))
        return "\n".join(lines)
    except ImportError:
        pass

    # Fallback: openpyxl
    from openpyxl import load_workbook
    wb = load_workbook(path, data_only=True)
    lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        lines.append(f"\n## Sheet: {sheet_name}")
        for row in ws.iter_rows(values_only=True):
            cells = [str(c) if c is not None else "" for c in row]
            if any(cells):
                lines.append("\t".join(cells))
    return "\n".join(lines)


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Extract text from Office files")
    parser.add_argument("input_file", help="Office file (.docx, .pptx, .xlsx)")
    parser.add_argument("--head", type=int, help="Show only first N lines", metavar="N")
    args = parser.parse_args()

    result = extract_text(args.input_file, head=args.head)
    print(result)
    sys.exit(1 if result.startswith("Error") else 0)
